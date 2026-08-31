import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_light_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Light UI Color Palette
    BG_CANVAS = RGBColor(248, 250, 252)       # #F8FAFC (Clean Light Slate)
    CARD_BG = RGBColor(255, 255, 255)         # #FFFFFF (Crisp White Card)
    CARD_BG_ALT = RGBColor(241, 245, 249)     # #F1F5F9 (Soft Slate)
    CARD_BORDER = RGBColor(226, 232, 240)     # #E2E8F0 (Subtle Slate Border)
    
    TEXT_NAVY = RGBColor(15, 39, 68)          # #0F2744 (Deep Slate Navy)
    TEXT_DARK = RGBColor(15, 23, 42)          # #0F172A (Dark Slate)
    TEXT_BODY = RGBColor(51, 65, 85)          # #334155 (Charcoal Slate)
    TEXT_MUTED = RGBColor(100, 116, 139)      # #64748B (Muted Slate)
    TEXT_WHITE = RGBColor(255, 255, 255)      # #FFFFFF (White)

    ACCENT_BLUE = RGBColor(29, 78, 216)       # #1D4ED8 (Royal Blue)
    ACCENT_SKY = RGBColor(2, 132, 199)        # #0284C7 (Electric Sky Blue)
    ACCENT_GREEN = RGBColor(5, 150, 105)      # #059669 (Emerald Green)
    ACCENT_MINT = RGBColor(16, 185, 129)      # #10B981 (Mint Green)
    ACCENT_GOLD = RGBColor(217, 119, 6)       # #D97706 (Amber Gold)
    ACCENT_YELLOW = RGBColor(245, 158, 11)    # #F59E0B (Golden Yellow)
    ACCENT_PURPLE = RGBColor(124, 58, 237)    # #7C3AED (Violet Purple)
    ACCENT_RED = RGBColor(225, 29, 72)        # #E11D48 (Crimson Rose)
    ACCENT_EMERALD_BG = RGBColor(240, 253, 244)# #F0FDF4 (Light Emerald Tint)

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_CANVAS
        bg.line.fill.background()
        return bg

    def add_slide_header(slide, category_tag, title_text, subtitle_text=None, tag_color=ACCENT_BLUE):
        # Category Tag Pill
        tag_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.42), Inches(3.8), Inches(0.35))
        tag_box.fill.solid()
        tag_box.fill.fore_color.rgb = CARD_BG
        tag_box.line.color.rgb = tag_color
        tag_box.line.width = Pt(1.5)
        tf = tag_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_tag.upper()
        p.font.name = "Segoe UI"
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = tag_color
        p.alignment = PP_ALIGN.CENTER

        # Title Textbox
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.733), Inches(0.85))
        tf_t = tb.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.name = "Segoe UI"
        p_t.font.size = Pt(24)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_NAVY

        if subtitle_text:
            p_s = tf_t.add_paragraph()
            p_s.text = subtitle_text
            p_s.font.name = "Segoe UI"
            p_s.font.size = Pt(12)
            p_s.font.color.rgb = TEXT_MUTED
            p_s.space_before = Pt(3)

    def add_slide_footer(slide, slide_num):
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.35))
        tf_f = footer_box.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = f"V.A.N.I - xAI | Bureau Of V.A.N.I-xAI (BoVxAi) | National Children's Science Congress (NCSC) | Slide {slide_num:02d}"
        p_f.font.name = "Segoe UI"
        p_f.font.size = Pt(10)
        p_f.font.color.rgb = TEXT_MUTED

    logo_path = os.path.join("VANI-B_IMGs", "Bureau_Of_VANI_xAI_Official_Logo.png")
    if not os.path.exists(logo_path):
        logo_path = "svg.png"

    # =========================================================================
    # SLIDE 1: Title & Executive Showcase (Pure Light UI)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1)

    # Master Light Hero Card
    hero_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.5), Inches(11.733), Inches(3.2))
    hero_card.fill.solid()
    hero_card.fill.fore_color.rgb = CARD_BG
    hero_card.line.color.rgb = ACCENT_BLUE
    hero_card.line.width = Pt(2)
    tf_hero = hero_card.text_frame
    tf_hero.word_wrap = True

    p0 = tf_hero.paragraphs[0]
    p0.text = "🏆 NATIONAL CHILDREN'S SCIENCE CONGRESS (NCSC) 2026-27"
    p0.font.name = "Segoe UI"
    p0.font.size = Pt(11.5)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_GOLD

    p1 = tf_hero.add_paragraph()
    p1.text = "V.A.N.I - xAI (BoVxAi)"
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_NAVY
    p1.space_before = Pt(2)

    p2 = tf_hero.add_paragraph()
    p2.text = "Vāṇī Adhyātmik Navīn Intellect"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(16)
    p2.font.color.rgb = ACCENT_GREEN
    p2.font.bold = True
    p2.space_before = Pt(2)

    p3 = tf_hero.add_paragraph()
    p3.text = "A Flagship Sovereign Innovation under the Bureau Of V.A.N.I-xAI (BoVxAi) | Next-Gen 100% Free AI Tutoring, Multimodal Socratic Mentorship & Autonomous Learning for Bharat"
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(12)
    p3.font.color.rgb = TEXT_BODY
    p3.space_before = Pt(4)

    # Logo picture on right
    if os.path.exists(logo_path):
        box1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(0.7), Inches(2.5), Inches(2.75))
        box1.fill.solid()
        box1.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box1.line.color.rgb = ACCENT_SKY
        box1.line.width = Pt(1.5)
        s1.shapes.add_picture(logo_path, Inches(10.0), Inches(0.85), width=Inches(2.1))

    # Two Cards Below
    c1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.9), Inches(5.7), Inches(2.95))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = ACCENT_BLUE
    c1.line.width = Pt(1.8)
    tf1 = c1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "👤 Innovator Profile"
    p.font.name = "Segoe UI"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    items1 = [
        ("Student Name:", "Dhruv Sagar"),
        ("Class & Stream:", "Class 10th (Secondary Section)"),
        ("Institution:", "PM SHRI KV NO.1 AFS Chakeri, Kanpur"),
        ("State / Region:", "Uttar Pradesh, India (KVS RO Varanasi)"),
        ("Student UID:", "1001049893")
    ]
    for k, v in items1:
        p_item = tf1.add_paragraph()
        p_item.text = f"• {k} {v}"
        p_item.font.name = "Segoe UI"
        p_item.font.size = Pt(11.5)
        p_item.font.color.rgb = TEXT_BODY
        p_item.space_before = Pt(3)

    c2 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(3.9), Inches(5.7), Inches(2.95))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = ACCENT_GREEN
    c2.line.width = Pt(1.8)
    tf2 = c2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "🎯 Core Innovation Highlights"
    p.font.name = "Segoe UI"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    items2 = [
        ("Parent Organization:", "Bureau Of V.A.N.I-xAI (BoVxAi) Research Cell"),
        ("Core Model:", "100% Free Sovereign AI (Zero Subscriptions & Zero Paywalls)"),
        ("Architecture:", "Asymmetric Cloud-Native Sub-User Engine (<50KB Client)"),
        ("Sustainability:", "100% Zero E-Waste & Zero Plastic Hardware"),
        ("Prototype Status:", "Fully Operational & Deployed Live on Web")
    ]
    for k, v in items2:
        p_item = tf2.add_paragraph()
        p_item.text = f"• {k} {v}"
        p_item.font.name = "Segoe UI"
        p_item.font.size = Pt(11.5)
        p_item.font.color.rgb = TEXT_BODY
        p_item.space_before = Pt(3)

    add_slide_footer(s1, 1)

    # =========================================================================
    # SLIDE 2: Innovator & Bureau Of V.A.N.I-xAI Profile
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    add_slide_header(s2, "Institutional Profile", "Innovator & Bureau Of V.A.N.I-xAI (BoVxAi)", "Grassroots technological leadership from PM SHRI KV No. 1 AFS Chakeri, Kanpur", ACCENT_BLUE)

    profiles = [
        ("The Student Innovator", "Dhruv Sagar (Class 10th)", [
            "Deep passion for Applied Machine Learning and Asymmetric Cloud Systems.",
            "Designed and hand-coded V.A.N.I-xAI from scratch to solve grassroots learning barriers.",
            "Specializes in lightweight edge client optimization for Tier-2 & Tier-3 students."
        ], ACCENT_BLUE),
        ("PM SHRI KV No. 1 Chakeri", "Excellence in Innovation", [
            "Selected under the prestigious PM SHRI Scheme by Government of India.",
            "Pioneering digital education, Atal Tinkering Labs, and grassroots science.",
            "Affiliated with CBSE and KVS Regional Office Varanasi / Lucknow."
        ], ACCENT_GREEN),
        ("Bureau Of V.A.N.I-xAI", "BoVxAi Sovereign Cell", [
            "Parent research cell for indigenous sovereign artificial intelligence.",
            "V.A.N.I AI serves as its flagship educational deployment.",
            "Committed to 100% zero paywalls, zero e-waste, and data sovereignty for India."
        ], ACCENT_GOLD)
    ]

    for idx, (title, subtitle, bullets, col) in enumerate(profiles):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(1.8)
        box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.75), Inches(5.0))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = col

        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = "Segoe UI"
        p_sub.font.size = Pt(12)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_DARK
        p_sub.space_before = Pt(3)

        for b in bullets:
            p_b = tf.add_paragraph()
            p_b.text = f"• {b}"
            p_b.font.name = "Segoe UI"
            p_b.font.size = Pt(11.5)
            p_b.font.color.rgb = TEXT_BODY
            p_b.space_before = Pt(8)

    add_slide_footer(s2, 2)

    # =========================================================================
    # SLIDE 3: Grassroots Problem Statement
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    add_slide_header(s3, "Grassroots Problem Identification", "The Digital & Pedagogical Divide in India", "Why generic commercial AI models fail 250M+ Indian school students", ACCENT_RED)

    problems = [
        ("1. Severe Economic Exclusivity & Paywalls", 
         "Commercial AI tools (ChatGPT Plus, Claude Pro) charge ₹1,500 to ₹2,500/month and mandate recurring credit card subscriptions. This locks out millions of talented underprivileged Indian students behind paywalls.", 
         ACCENT_RED),
        ("2. Hardware & Low-Bandwidth Barrier", 
         "Advanced generative models demand high-end GPUs, latest smartphones, and high-speed broadband. Underprivileged students using ₹5,000 smartphones and 2G/3G mobile data suffer crashes and freezes.", 
         ACCENT_GOLD),
        ("3. The 'Cheat-Engine' Flaw", 
         "Mainstream LLMs generate direct solutions without explanations. This encourages rote copy-pasting and academic cheating rather than cognitive reasoning, critical inquiry, or conceptual mastery.", 
         ACCENT_PURPLE),
        ("4. Lack of Localized & Safe Guardrails", 
         "Foreign AI platforms lack alignment with Indian NCERT/CBSE curriculum, Vernacular voice capabilities in Hindi/dialects, and youth-safe ethical guardrails tailored for Indian schools.", 
         ACCENT_BLUE)
    ]

    for idx, (head, desc, col) in enumerate(problems):
        row = idx // 2
        col_idx = idx % 2
        left = Inches(0.8 + col_idx * 5.95)
        top = Inches(1.8 + row * 2.55)
        box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.35))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = head
        p.font.name = "Segoe UI"
        p.font.size = Pt(14.5)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = TEXT_BODY
        p2.space_before = Pt(6)

    add_slide_footer(s3, 3)

    # =========================================================================
    # SLIDE 4: The Proposed Innovation
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    add_slide_header(s4, "The Proposed Innovation", "V.A.N.I-xAI: Sovereign AI Mentor for Bharat", "Democratizing cloud intelligence through indigenous asymmetric architecture", ACCENT_GREEN)

    solutions = [
        ("1. <50KB Asymmetric Core", 
         "99% of compute is shifted to cloud clusters. The web client payload is <50KB, allowing instant sub-second boot times even on legacy ₹5,000 mobile phones and slow 2G/3G connections.",
         ACCENT_BLUE),
        ("2. Socratic Mentorship Engine", 
         "Instead of regurgitating direct answers, the model acts as an empathetic digital mentor—breaking complex STEM concepts into interactive hints, guided questioning, and conceptual scaffolding.",
         ACCENT_GREEN),
        ("3. BoV Multimodal Voice Hub", 
         "Integrated speech synthesis and voice processing (voice_agent.py) enables real-time auditory explanations, assisting students with reading difficulties or vernacular preferences.",
         ACCENT_PURPLE),
        ("4. 100% Free Sovereign Access", 
         "Removed all premium subscription plans and paywalls. Every capability is 100% free and permanently unlocked for all students with seamless Google OAuth and mathematical UID verification.",
         ACCENT_GOLD)
    ]

    for idx, (title, desc, col) in enumerate(solutions):
        row = idx // 2
        col_idx = idx % 2
        left = Inches(0.8 + col_idx * 5.95)
        top = Inches(1.8 + row * 2.55)
        box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.35))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(14.5)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = TEXT_BODY
        p2.space_before = Pt(6)

    add_slide_footer(s4, 4)

    # =========================================================================
    # SLIDE 5: Bureau of V.A.N.I-xAI (BoVxAi) Novel IP
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5)
    add_slide_header(s5, "Scientific Novelty & IP", "Bureau Of V.A.N.I-xAI (BoVxAi) 4 Novel Pillars", "Handcrafted algorithmic systems engineered from scratch for the V.A.N.I ecosystem", ACCENT_GOLD)

    pillars = [
        ("1. Mathematical UID Token", 
         "Algorithmically generates permanent tokens (e.g. V.A.N.I-xAI-DHRU-3478) upon Google OAuth sign-in, locking persistent memory and multi-device workspaces without hardware lock-in.", 
         ACCENT_BLUE),
        ("2. 3-Layer Security Check", 
         "Proprietary algorithmic authorization validated against 3 rigorous checks: Database Existence, Active Student Token validation, and strict UID sandbox isolation.", 
         ACCENT_GREEN),
        ("3. InAixVAi Dashboard", 
         "A bespoke educator/admin dashboard enabling real-time multi-user telemetry, classroom activity monitoring, query load oversight, and platform integrity management.", 
         ACCENT_GOLD),
        ("4. Autonomous Memory Engine", 
         "An autonomous background self-healing engine managing real-time cross-device session synchronization, SQLite memory persistence (vani_memory.db), and load-balancing.", 
         ACCENT_PURPLE)
    ]

    for idx, (title, desc, col) in enumerate(pillars):
        left = Inches(0.8 + idx * 2.95)
        top = Inches(1.8)
        box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(5.0))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_BODY
        p2.space_before = Pt(8)

    add_slide_footer(s5, 5)

    # =========================================================================
    # SLIDE 6: Multi-Agent Intelligence & Code Modules
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6)
    add_slide_header(s6, "Applied AI Capabilities", "Multi-Agent Intelligence & Backend Modules", "State-of-the-art intelligent modules integrated into the local & cloud backend", ACCENT_BLUE)

    agent_features = [
        ("🧠 LLM Router & SQLite Memory (llm_router.py)", 
         "Dynamic multi-model fallback across vision and text models. Backed by vani_memory.db for persistent conversational memory, contextual recall, and fast semantic search.", 
         ACCENT_BLUE),
        ("🎙️ Multimodal Voice Hub (voice_agent.py)", 
         "Powered by edge TTS and speech processing. Converts complex scientific derivations and guided explanations into natural audio speech for inclusive hands-free learning.", 
         ACCENT_GREEN),
        ("🤖 Vision-Guided Agentic Loop (agentic_loop.py)", 
         "Utilizes visual reasoning to analyze user screen states, formulas, and educational diagrams, transforming V.A.N.I into an interactive digital lab assistant.", 
         ACCENT_GOLD),
        ("⚡ Proactive System Agent (autonomous_agent.py)", 
         "Monitors real-time battery status, initiates eco-friendly power-saving protocols, and schedules automated learning routines for student productivity.", 
         ACCENT_PURPLE)
    ]

    for idx, (title, desc, col) in enumerate(agent_features):
        row = idx // 2
        col_idx = idx % 2
        left = Inches(0.8 + col_idx * 5.95)
        top = Inches(1.8 + row * 2.55)
        box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.35))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = TEXT_BODY
        p2.space_before = Pt(6)

    add_slide_footer(s6, 6)

    # =========================================================================
    # SLIDE 7: 5-Step Operational Flow
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7)
    add_slide_header(s7, "End-to-End Pipeline", "5-Step Zero-Friction Operational Pipeline", "Frictionless workflow from Google authentication to persistent concept mastery", ACCENT_GREEN)

    steps = [
        ("Step 1", "Auth & UID Gen", "Student signs in securely via Google OAuth. System computes and assigns permanent V.A.N.I UID.", ACCENT_BLUE),
        ("Step 2", "Instant Free Access", "Student enters workspace with 100% unlocked features, zero subscription barriers, zero fees.", ACCENT_GREEN),
        ("Step 3", "Security & Swarm", "Binds UID with session tokens, enabling instant cross-device pairing and encrypted state sync.", ACCENT_GOLD),
        ("Step 4", "Socratic Mentoring", "Queries route through LLM router, memory database, and Socratic filters for guided tutoring.", ACCENT_PURPLE),
        ("Step 5", "Continuous Learning", "Autonomous memory engine continuously stores learning milestones for lifelong growth.", ACCENT_BLUE)
    ]

    for idx, (step_num, title, desc, col) in enumerate(steps):
        left = Inches(0.8 + idx * 2.38)
        top = Inches(1.9)
        box = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.25), Inches(4.8))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        
        tf = box.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = step_num.upper()
        p0.font.name = "Segoe UI"
        p0.font.size = Pt(11)
        p0.font.bold = True
        p0.font.color.rgb = col
        
        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_DARK
        p1.space_before = Pt(4)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_BODY
        p2.space_before = Pt(8)

    add_slide_footer(s7, 7)

    # =========================================================================
    # SLIDE 8: Comparative Advantage Matrix
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8)
    add_slide_header(s8, "Competitive Benchmarking", "V.A.N.I-xAI vs Conventional Platforms", "Engineered specifically for maximum grassroots accessibility & pedagogical integrity", ACCENT_GOLD)

    rows = 6
    cols = 4
    left = Inches(0.8)
    top = Inches(1.8)
    width = Inches(11.733)
    height = Inches(4.9)

    table_shape = s8.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.columns[0].width = Inches(3.1)
    table.columns[1].width = Inches(2.85)
    table.columns[2].width = Inches(2.85)
    table.columns[3].width = Inches(2.933)

    headers = ["Evaluation Metric", "Commercial AI (ChatGPT/Claude)", "Traditional EdTech Apps", "V.A.N.I - xAI (BoVxAi)"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEXT_NAVY if i != 3 else ACCENT_GREEN
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.name = "Segoe UI"
        p.font.bold = True
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_WHITE

    data = [
        ("Cost for Students", "₹1,500 - ₹2,500/month (Expensive)", "₹5,000 - ₹25,000/year (High cost)", "100% Free Lifetime (Zero Paywalls)"),
        ("Hardware / Device Barrier", "Requires modern smartphone & RAM", "Requires 100MB+ app download", "Runs on any browser / <50KB payload"),
        ("Pedagogical Approach", "Direct answer generation (cheat)", "Static recorded video lectures", "Socratic guidance & Concept scaffolding"),
        ("Payment & Access Barrier", "Mandates International Credit Cards", "Closed recurring subscriptions", "Zero Paywalls / Instant Free Access"),
        ("Zero E-Waste & Ecology", "Heavy local battery & GPU drain", "Proprietary locked tablet hardware", "100% Cloud-Native (Zero E-Waste)")
    ]

    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            if c_idx == 3:
                cell.fill.fore_color.rgb = ACCENT_EMERALD_BG
            else:
                cell.fill.fore_color.rgb = CARD_BG if r_idx % 2 == 0 else CARD_BG_ALT
                
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Segoe UI"
            p.font.size = Pt(10.5)
            if c_idx == 3:
                p.font.bold = True
                p.font.color.rgb = ACCENT_GREEN
            elif c_idx == 0:
                p.font.bold = True
                p.font.color.rgb = TEXT_DARK
            else:
                p.font.color.rgb = TEXT_BODY

    add_slide_footer(s8, 8)

    # =========================================================================
    # SLIDE 9: Zero E-Waste & Environmental Sustainability
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9)
    add_slide_header(s9, "Environmental Sustainability", "100% Zero E-Waste & Eco-Friendly Architecture", "A sustainable software innovation aligned with National Green Missions", ACCENT_GREEN)

    c1 = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5.0))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = ACCENT_GREEN
    c1.line.width = Pt(1.8)
    tf1 = c1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "🌱 Zero Physical Manufacturing"
    p.font.name = "Segoe UI"
    p.font.size = Pt(14.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    
    points1 = [
        "100% software-driven: Requires zero physical plastic casing, toxic lithium batteries, or PCB manufacturing.",
        "Completely eliminates electronic scrap at the manufacturing root.",
        "Drastically reduces carbon emissions through optimized cloud server inference over heavy local processing."
    ]
    for pt in points1:
        p = tf1.add_paragraph()
        p.text = f"• {pt}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_BODY
        p.space_before = Pt(8)

    c2 = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5.0))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = ACCENT_BLUE
    c2.line.width = Pt(1.8)
    tf2 = c2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "📱 Extends Old Hardware Lifespan"
    p.font.name = "Segoe UI"
    p.font.size = Pt(14.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    points2 = [
        "Repurposes existing school computer labs and parents' low-cost old smartphones.",
        "Prevents premature disposal of functional electronics by keeping client footprint <50KB.",
        "Runs directly inside standard web browsers with zero installation requirements."
    ]
    for pt in points2:
        p = tf2.add_paragraph()
        p.text = f"• {pt}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_BODY
        p.space_before = Pt(8)

    c3 = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0))
    c3.fill.solid()
    c3.fill.fore_color.rgb = CARD_BG
    c3.line.color.rgb = ACCENT_GOLD
    c3.line.width = Pt(1.8)
    tf3 = c3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "📄 100% Paperless Ecosystem"
    p.font.name = "Segoe UI"
    p.font.size = Pt(14.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    
    points3 = [
        "Eliminates physical question banks, printed worksheets, and test guides.",
        "All learning sessions, notes, and milestones are preserved digitally in persistent memory.",
        "Supports India's Green Technology and Net-Zero Digital Education initiatives."
    ]
    for pt in points3:
        p = tf3.add_paragraph()
        p.text = f"• {pt}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_BODY
        p.space_before = Pt(8)

    add_slide_footer(s9, 9)

    # =========================================================================
    # SLIDE 10: Societal Impact & DBT ₹10,000 Fund Budget
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10)
    add_slide_header(s10, "Societal Value & Scalability", "Grassroots Value & Scaling Plan", "Resource optimization and maximum educational impact for Bharat", ACCENT_GOLD)

    card_left = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = CARD_BG
    card_left.line.color.rgb = ACCENT_BLUE
    card_left.line.width = Pt(1.8)
    tf_l = card_left.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "🤝 Grassroots Societal Value"
    p.font.name = "Segoe UI"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    soc_points = [
        "Democratizing Tutoring: Brings personalized 1-on-1 AI guidance to Tier-2, Tier-3 & village learners at ₹0 cost.",
        "Empowers First-Gen Learners: Deep conceptual clarity in Science, Math, and Computer Programming.",
        "NEP 2020 Alignment: Fosters experiential and inquiry-driven learning over rote memorization.",
        "Inclusivity: Voice engine assists students with reading difficulties or vernacular preferences."
    ]
    for pt in soc_points:
        p = tf_l.add_paragraph()
        p.text = f"• {pt}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_BODY
        p.space_before = Pt(8)

    card_right = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(1.8), Inches(5.7), Inches(5.0))
    card_right.fill.solid()
    card_right.fill.fore_color.rgb = CARD_BG
    card_right.line.color.rgb = ACCENT_GOLD
    card_right.line.width = Pt(1.8)
    tf_r = card_right.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "💰 DST ₹10,000 DBT Fund Utilization Plan"
    p.font.name = "Segoe UI"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD

    dbt_points = [
        ("₹5,000 (50%)", "High-speed Cloud API Credits: Empowers 100,000+ student educational queries at zero cost to families."),
        ("₹2,500 (25%)", "Serverless Database & Edge CDN: High-speed low-latency cloud hosting across Indian regional nodes."),
        ("₹1,500 (15%)", "Vernacular Language Tokenization: Optimization for Hindi and regional Indian dialects."),
        ("₹1,000 (10%)", "Security & Domain Infrastructure: Enterprise SSL certificates & school cluster deployment.")
    ]
    for k, v in dbt_points:
        p = tf_r.add_paragraph()
        p.text = f"• {k}: {v}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_BODY
        p.space_before = Pt(8)

    add_slide_footer(s10, 10)

    # =========================================================================
    # SLIDE 11: Live Working Prototype Readiness
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_background(s11)
    add_slide_header(s11, "Working Prototype Demonstration", "Fully Operational & Deployed Features", "Tested and verified codebase ready for school & district implementation", ACCENT_GREEN)

    features = [
        ("🔐 InAixVAi Command Panel", 
         "A live administrative panel enabling teachers and admins to track authenticated sub-users, monitor query loads, manage classroom sessions, and safeguard student data in real-time.", 
         ACCENT_BLUE),
        ("🧠 LLM Router & Semantic Context", 
         "Equipped with multi-model fallback, local memory database (vani_memory.db), and Socratic prompt pipelines to ensure accurate, hallucination-resistant educational mentorship.", 
         ACCENT_GREEN),
        ("🎙️ Multimodal Voice Tutor (BoV)", 
         "Integrated voice synthesizer (voice_agent.py) converting written solutions to spoken words, enabling auditory learners and visually impaired students to engage effortlessly.", 
         ACCENT_GOLD),
        ("📱 Universal Device Compatibility", 
         "Zero-install, browser-accessible interface tested across low-end Android phones, iPads, Windows PCs, and Linux school lab workstations with <1s load latency.", 
         ACCENT_PURPLE)
    ]

    for idx, (title, desc, col) in enumerate(features):
        row = idx // 2
        col_idx = idx % 2
        left = Inches(0.8 + col_idx * 5.95)
        top = Inches(1.8 + row * 2.55)
        box = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.35))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = TEXT_BODY
        p2.space_before = Pt(6)

    add_slide_footer(s11, 11)

    # =========================================================================
    # SLIDE 12: Vision 2030 Roadmap & Grand Conclusion (Pure Light UI)
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_background(s12)
    add_slide_header(s12, "Future Scope & Vision 2030", "Roadmap: From School Innovation to National Utility", "Structured scale-up strategy for NCSC & National Digital Education Missions", ACCENT_GOLD)

    phases = [
        ("Phase 1: School Pilot (2026)", 
         "Ready-to-Deploy Prototype", 
         "• Sovereign sub-user architecture deployed\n• Socratic AI engine with zero-cost onboarding\n• InAixVAi teacher administrative dashboard\n• Pilot launch at PM SHRI KV No. 1 Chakeri", 
         ACCENT_BLUE),
        ("Phase 2: District & State Level", 
         "Vernacular & Voice Scaling", 
         "• Integration of 10+ Indian regional languages\n• Bi-directional voice conversation engine\n• Offline mesh/edge caching for remote village schools\n• State-wide KV & Government school rollout", 
         ACCENT_GREEN),
        ("Phase 3: National Level", 
         "Public Sovereign AI Ecosystem", 
         "• Potential integration with DIKSHA & PM eVidya\n• Interactive AI virtual science & coding lab assistant\n• Complete Atmanirbhar Bharat AI mentoring stack", 
         ACCENT_GOLD)
    ]

    for idx, (phase, subtitle, bullets, col) in enumerate(phases):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(1.8)
        box = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.75), Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = phase
        p.font.name = "Segoe UI"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = "Segoe UI"
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = TEXT_DARK
        p_sub.font.bold = True
        p_sub.space_before = Pt(3)

        p_b = tf.add_paragraph()
        p_b.text = bullets
        p_b.font.name = "Segoe UI"
        p_b.font.size = Pt(10.5)
        p_b.font.color.rgb = TEXT_BODY
        p_b.space_before = Pt(6)

    # Bottom Light UI Grand Banner
    banner = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.5), Inches(11.733), Inches(1.3))
    banner.fill.solid()
    banner.fill.fore_color.rgb = ACCENT_EMERALD_BG
    banner.line.color.rgb = ACCENT_GREEN
    banner.line.width = Pt(1.5)
    tf_ban = banner.text_frame
    tf_ban.word_wrap = True
    
    p_ban = tf_ban.paragraphs[0]
    p_ban.text = "🏆 V.A.N.I-xAI &bull; Bureau Of V.A.N.I-xAI (BoVxAi) | Atmanirbhar Bharat AI"
    p_ban.font.name = "Segoe UI"
    p_ban.font.size = Pt(13)
    p_ban.font.bold = True
    p_ban.font.color.rgb = ACCENT_GREEN
    p_ban.alignment = PP_ALIGN.CENTER

    p_ban2 = tf_ban.add_paragraph()
    p_ban2.text = "Thank you, Respected Jury Members & NCSC Mentors | Ready to Empower Every Indian Student"
    p_ban2.font.name = "Segoe UI"
    p_ban2.font.size = Pt(11.5)
    p_ban2.font.bold = True
    p_ban2.font.color.rgb = TEXT_NAVY
    p_ban2.alignment = PP_ALIGN.CENTER
    p_ban2.space_before = Pt(3)

    add_slide_footer(s12, 12)

    output_filename = "VANI_xAI_NCSC_2026_Light.pptx"
    prs.save(output_filename)
    print(f"Successfully generated 100% Light UI PPTX: {output_filename}")
    return output_filename

if __name__ == "__main__":
    build_light_presentation()
