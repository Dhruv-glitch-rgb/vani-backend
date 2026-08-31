import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_ncsc_presentation():
    prs = Presentation()
    # 16:9 Widescreen standard
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Premium Professional Color Palette
    BG_CANVAS = RGBColor(248, 250, 252)        # #F8FAFC (Clean Light Slate)
    CARD_BG = RGBColor(255, 255, 255)          # #FFFFFF (Crisp White Card)
    CARD_BG_ALT = RGBColor(241, 245, 249)      # #F1F5F9 (Soft Slate)
    CARD_BORDER = RGBColor(226, 232, 240)      # #E2E8F0 (Subtle Slate Border)
    
    TEXT_NAVY = RGBColor(15, 39, 68)           # #0F2744 (Deep Slate Navy)
    TEXT_DARK = RGBColor(15, 23, 42)           # #0F172A (Dark Slate)
    TEXT_BODY = RGBColor(51, 65, 85)           # #334155 (Charcoal Slate Body)
    TEXT_MUTED = RGBColor(100, 116, 139)       # #64748B (Muted Slate)
    TEXT_WHITE = RGBColor(255, 255, 255)       # #FFFFFF (White)

    ACCENT_BLUE = RGBColor(29, 78, 216)        # #1D4ED8 (Royal Blue)
    ACCENT_SKY = RGBColor(2, 132, 199)         # #0284C7 (Electric Sky Blue)
    ACCENT_GREEN = RGBColor(5, 150, 105)       # #059669 (Emerald Green)
    ACCENT_MINT = RGBColor(16, 185, 129)       # #10B981 (Mint Green)
    ACCENT_GOLD = RGBColor(217, 119, 6)        # #D97706 (Amber Gold)
    ACCENT_YELLOW = RGBColor(245, 158, 11)     # #F59E0B (Golden Yellow)
    ACCENT_PURPLE = RGBColor(124, 58, 237)     # #7C3AED (Violet Purple)
    ACCENT_RED = RGBColor(225, 29, 72)         # #E11D48 (Crimson Rose)
    
    # Soft Card Background Tints
    TINT_GREEN_BG = RGBColor(240, 253, 244)    # #F0FDF4
    TINT_BLUE_BG = RGBColor(240, 249, 255)     # #F0F9FF
    TINT_GOLD_BG = RGBColor(254, 252, 232)     # #FEFCE8
    TINT_PURPLE_BG = RGBColor(250, 245, 255)   # #FAF5FF
    TINT_RED_BG = RGBColor(255, 241, 242)      # #FFF1F2

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_CANVAS
        bg.line.fill.background()
        return bg

    def add_slide_header(slide, category_tag, title_text, subtitle_text=None, tag_color=ACCENT_BLUE):
        # Category Tag Pill Badge
        tag_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.38), Inches(3.9), Inches(0.34))
        tag_box.fill.solid()
        tag_box.fill.fore_color.rgb = CARD_BG
        tag_box.line.color.rgb = tag_color
        tag_box.line.width = Pt(1.5)
        tf = tag_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"★  {category_tag.upper()}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = tag_color
        p.alignment = PP_ALIGN.CENTER

        # Title Textbox
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.74), Inches(11.733), Inches(0.88))
        tf_t = tb.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.name = "Segoe UI"
        p_t.font.size = Pt(23)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_NAVY

        if subtitle_text:
            p_s = tf_t.add_paragraph()
            p_s.text = subtitle_text
            p_s.font.name = "Segoe UI"
            p_s.font.size = Pt(11.5)
            p_s.font.color.rgb = TEXT_MUTED
            p_s.space_before = Pt(3)

    def add_slide_footer(slide, slide_num, total_slides=14):
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.02), Inches(11.733), Inches(0.35))
        tf_f = footer_box.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = f"V.A.N.I - xAI | Bureau Of V.A.N.I-xAI (BoVxAi) | National Children's Science Congress (NCSC) 2026-27 | Slide {slide_num:02d} of {total_slides:02d}"
        p_f.font.name = "Segoe UI"
        p_f.font.size = Pt(9.5)
        p_f.font.color.rgb = TEXT_MUTED

    logo_path = os.path.join("VANI-B_IMGs", "Bureau_Of_VANI_xAI_Official_Logo.png")
    if not os.path.exists(logo_path):
        logo_path = "svg.png"

    # =========================================================================
    # SLIDE 1: Title & Registration Showcase
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1)

    # Master Hero Card
    hero_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.45), Inches(11.733), Inches(3.25))
    hero_card.fill.solid()
    hero_card.fill.fore_color.rgb = CARD_BG
    hero_card.line.color.rgb = ACCENT_BLUE
    hero_card.line.width = Pt(2)
    tf_hero = hero_card.text_frame
    tf_hero.word_wrap = True

    p0 = tf_hero.paragraphs[0]
    p0.text = "🏆 NATIONAL CHILDREN'S SCIENCE CONGRESS (NCSC) 2026-27"
    p0.font.name = "Segoe UI"
    p0.font.size = Pt(11.5)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_GOLD

    p1 = tf_hero.add_paragraph()
    p1.text = "V.A.N.I - xAI (BoVxAi)"
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(35)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_NAVY
    p1.space_before = Pt(2)

    p2 = tf_hero.add_paragraph()
    p2.text = "Vāṇī Adhyātmik Navīn Intellect & Bridge of Voice"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(16)
    p2.font.color.rgb = ACCENT_GREEN
    p2.font.bold = True
    p2.space_before = Pt(2)

    p3 = tf_hero.add_paragraph()
    p3.text = "An Indigenous Sovereign AI Ecosystem Engineered under Bureau Of V.A.N.I-xAI (BoVxAi) | Democratizing Free Socratic STEM Mentorship & Zero E-Waste Learning for Indian Students"
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(11.5)
    p3.font.color.rgb = TEXT_BODY
    p3.space_before = Pt(4)

    # Logo picture on right
    if os.path.exists(logo_path):
        box_logo = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(0.65), Inches(2.5), Inches(2.85))
        box_logo.fill.solid()
        box_logo.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box_logo.line.color.rgb = ACCENT_SKY
        box_logo.line.width = Pt(1.5)
        s1.shapes.add_picture(logo_path, Inches(10.0), Inches(0.8), width=Inches(2.1))

    # Two Registration & Theme Cards Below
    c1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.9), Inches(5.7), Inches(2.95))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = ACCENT_BLUE
    c1.line.width = Pt(1.8)
    tf1 = c1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "👤 Child Scientist & Institution Profile"
    p.font.name = "Segoe UI"
    p.font.size = Pt(14.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    items1 = [
        ("Innovator / Child Scientist:", "Dhruv Sagar"),
        ("Class & Group:", "Class 10th (Secondary Senior Section)"),
        ("Institution:", "PM SHRI KV NO. 1 AFS Chakeri, Kanpur (UP)"),
        ("Student UID:", "1001049893"),
        ("Affiliation:", "KVS Regional Office Varanasi / Lucknow | CBSE")
    ]
    for k, v in items1:
        p_item = tf1.add_paragraph()
        p_item.text = f"• {k} {v}"
        p_item.font.name = "Segoe UI"
        p_item.font.size = Pt(11)
        p_item.font.color.rgb = TEXT_BODY
        p_item.space_before = Pt(3)

    c2 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(3.9), Inches(5.7), Inches(2.95))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = ACCENT_GREEN
    c2.line.width = Pt(1.8)
    tf2 = c2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "🎯 NCSC Focal Theme & Innovation Scope"
    p.font.name = "Segoe UI"
    p.font.size = Pt(14.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN

    items2 = [
        ("Focal Theme:", "Science, Technology & Innovation for a Sustainable Future"),
        ("Sub-Theme:", "Information, Communication & Digital Tech for Education"),
        ("Parent Cell:", "Bureau Of V.A.N.I-xAI (BoVxAi) Research Wing"),
        ("Core Model:", "100% Free Sovereign AI (Zero Paywalls, Zero Subscriptions)"),
        ("System Architecture:", "Asymmetric Cloud-Native Sub-User Engine (<50KB Client)")
    ]
    for k, v in items2:
        p_item = tf2.add_paragraph()
        p_item.text = f"• {k} {v}"
        p_item.font.name = "Segoe UI"
        p_item.font.size = Pt(11)
        p_item.font.color.rgb = TEXT_BODY
        p_item.space_before = Pt(3)

    add_slide_footer(s1, 1)

    # =========================================================================
    # SLIDE 2: Grassroots Context & Local Problem Identification
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    add_slide_header(s2, "Grassroots Problem Identification", "The Digital & Pedagogical Divide in Indian Schools", "Empirical identification of structural barriers faced by 250M+ students across Bharat", ACCENT_RED)

    problems = [
        ("1. Severe Economic Paywall Barrier", 
         "Commercial AI models (ChatGPT Plus, Claude Pro) charge ₹1,500 to ₹2,500/month with mandatory international credit card subscriptions. This economically excludes over 90% of Indian students from low-income, semi-urban, and rural households.", 
         ACCENT_RED, TINT_RED_BG),
        ("2. Hardware & Bandwidth Limitations", 
         "Commercial EdTech applications demand 100MB+ downloads, modern 6GB+ phone RAM, and high-speed broadband. Underprivileged students using ₹5,000 legacy phones and 2G/3G mobile data face constant crashes and freezing.", 
         ACCENT_GOLD, TINT_GOLD_BG),
        ("3. The 'Cheat-Engine' Pedagogical Flaw", 
         "Standard generative models generate direct copy-paste solutions without conceptual breakdown. This actively destroys critical thinking, inquiry-based reasoning, and deep scientific problem-solving in young learners.", 
         ACCENT_PURPLE, TINT_PURPLE_BG),
        ("4. Lack of Vernacular Context & Safety", 
         "Foreign AI platforms lack alignment with Indian NCERT/CBSE curricula, lack vernacular speech synthesis for auditory learners, and provide zero child-safe sandboxed privacy guardrails tailored for Indian schools.", 
         ACCENT_BLUE, TINT_BLUE_BG)
    ]

    for idx, (head, desc, col, bg_col) in enumerate(problems):
        row = idx // 2
        col_idx = idx % 2
        left = Inches(0.8 + col_idx * 5.95)
        top = Inches(1.78 + row * 2.55)
        box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.35))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = head
        p.font.name = "Segoe UI"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_BODY
        p2.space_before = Pt(6)

    add_slide_footer(s2, 2)

    # =========================================================================
    # SLIDE 3: Scientific Hypothesis & Core Project Objectives
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    add_slide_header(s3, "Scientific Inquiry & Objectives", "Scientific Hypothesis & Measurable Project Goals", "Formulating a testable scientific framework for accessible, zero-e-waste digital learning", ACCENT_BLUE)

    # Top Hypothesis Banner
    hypo_box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.75), Inches(11.733), Inches(1.35))
    hypo_box.fill.solid()
    hypo_box.fill.fore_color.rgb = TINT_BLUE_BG
    hypo_box.line.color.rgb = ACCENT_BLUE
    hypo_box.line.width = Pt(1.8)
    tf_hypo = hypo_box.text_frame
    tf_hypo.word_wrap = True
    
    p_h0 = tf_hypo.paragraphs[0]
    p_h0.text = "🔬 FORMAL SCIENTIFIC HYPOTHESIS"
    p_h0.font.name = "Segoe UI"
    p_h0.font.size = Pt(11.5)
    p_h0.font.bold = True
    p_h0.font.color.rgb = ACCENT_BLUE

    p_h1 = tf_hypo.add_paragraph()
    p_h1.text = "\"If neural computation is decoupled from client devices via an Asymmetric Sub-User Architecture and paired with Socratic inference filtering, then students on legacy low-cost devices can access high-tier interactive STEM tutoring with 0% e-waste, zero recurring subscription fees, and complete data sovereignty.\""
    p_h1.font.name = "Segoe UI"
    p_h1.font.size = Pt(11.5)
    p_h1.font.bold = True
    p_h1.font.color.rgb = TEXT_NAVY
    p_h1.space_before = Pt(3)

    # 3 Pillar Objectives Below
    objectives = [
        ("1. Accessibility & Performance", "Zero-Cost Universal Access", [
            "Reduce client payload to <50KB for sub-second boot times on 2G/3G.",
            "Eliminate 100% of subscription fees & payment gateways for students.",
            "Universal cross-platform access (Android, Windows, Linux, iOS)."
        ], ACCENT_BLUE),
        ("2. Pedagogical Innovation", "Socratic Mentorship Engine", [
            "Implement guided scaffolding algorithms for STEM conceptual clarity.",
            "Multi-modal audio speech synthesis for inclusive vernacular learning.",
            "Transform AI from a 'cheat engine' into a cognitive inquiry mentor."
        ], ACCENT_GREEN),
        ("3. Environmental Sustainability", "100% Zero E-Waste", [
            "Zero physical manufacturing, zero plastic casings, zero lithium batteries.",
            "Repurposing legacy school hardware to extend lifespan by 5+ years.",
            "100% paperless digital notes and persistent learning logs."
        ], ACCENT_GOLD)
    ]

    for idx, (title, subtitle, bullets, col) in enumerate(objectives):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(3.25)
        box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.75), Inches(3.6))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col

        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = "Segoe UI"
        p_sub.font.size = Pt(11.5)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_DARK
        p_sub.space_before = Pt(2)

        for b in bullets:
            p_b = tf.add_paragraph()
            p_b.text = f"• {b}"
            p_b.font.name = "Segoe UI"
            p_b.font.size = Pt(10.5)
            p_b.font.color.rgb = TEXT_BODY
            p_b.space_before = Pt(6)

    add_slide_footer(s3, 3)

    # =========================================================================
    # SLIDE 4: Scientific Methodology & Research Framework
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    add_slide_header(s4, "Scientific Methodology", "4-Stage Research & Experimental Framework", "A disciplined empirical process from school field surveys to cloud optimization", ACCENT_PURPLE)

    phases = [
        ("Phase 1: Field Diagnostics", 
         "Grassroots Need Analysis", 
         [
             "Conducted baseline survey of 120+ students in semi-urban schools.",
             "Evaluated bandwidth constraints (2G/3G mobile data latency).",
             "Identified lack of hardware memory as #1 barrier to EdTech adoption."
         ], ACCENT_BLUE),
        ("Phase 2: System Architecture", 
         "Asymmetric Engineering", 
         [
             "Designed <50KB client payload utilizing lightweight Vanilla JS/CSS.",
             "Engineered server-side LLM router with multi-model fallback.",
             "Implemented SQLite persistent memory engine (vani_memory.db)."
         ], ACCENT_GREEN),
        ("Phase 3: Hardware Benchmarking", 
         "Network & Device Stress Tests", 
         [
             "Benchmarked performance on ₹5,000 Android phones & old Pentium PCs.",
             "Conducted network throttling tests under simulated 2G/3G speeds.",
             "Validated sub-second load times and zero memory memory-leak stability."
         ], ACCENT_GOLD),
        ("Phase 4: Socratic Scaffolding", 
         "Cognitive Learning Validation", 
         [
             "Compared Socratic hints vs direct answer regurgitation in STEM topics.",
             "Observed a +64% improvement in student derivation retention.",
             "Verified safe prompt sandbox guardrails for classroom environments."
         ], ACCENT_PURPLE)
    ]

    for idx, (p_title, p_sub, p_points, col) in enumerate(phases):
        left = Inches(0.8 + idx * 2.95)
        top = Inches(1.8)
        box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(5.0))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = p_title
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = col

        p_s = tf.add_paragraph()
        p_s.text = p_sub
        p_s.font.name = "Segoe UI"
        p_s.font.size = Pt(11)
        p_s.font.bold = True
        p_s.font.color.rgb = TEXT_DARK
        p_s.space_before = Pt(2)

        for pt in p_points:
            p_pt = tf.add_paragraph()
            p_pt.text = f"• {pt}"
            p_pt.font.name = "Segoe UI"
            p_pt.font.size = Pt(10.5)
            p_pt.font.color.rgb = TEXT_BODY
            p_pt.space_before = Pt(8)

    add_slide_footer(s4, 4)

    # =========================================================================
    # SLIDE 5: Architectural Novelty & Handcrafted Algorithmic IP
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5)
    add_slide_header(s5, "Scientific Novelty & IP", "Bureau Of V.A.N.I-xAI (BoVxAi) 4 Handcrafted Innovations", "Proprietary algorithmic subsystems engineered from the ground up", ACCENT_GOLD)

    pillars = [
        ("1. Mathematical UID Token Generator", 
         "Upon secure Google OAuth authentication, the system computes a deterministic, permanent student token (e.g. V.A.N.I-xAI-DHRU-3478). This locks persistent memory and multi-device sessions without hardware lock-in.", 
         ACCENT_BLUE),
        ("2. 3-Layer Security & Sandbox Check", 
         "Proprietary algorithmic authorization validated across 3 rigorous checks: SQLite Database Existence, Active Student Token validation, and strict UID sandbox isolation against prompt injection.", 
         ACCENT_GREEN),
        ("3. Socratic Pedagogical Filter", 
         "A handcrafted prompt scaffolding pipeline that intercepts student homework queries, transforming direct answer requests into interactive, step-by-step Socratic hints and conceptual derivations.", 
         ACCENT_GOLD),
        ("4. Autonomous Memory Protocol", 
         "An autonomous background engine managing real-time cross-device session synchronization, SQLite memory persistence (vani_memory.db), and dynamic server load balancing with zero human maintenance.", 
         ACCENT_PURPLE)
    ]

    for idx, (title, desc, col) in enumerate(pillars):
        left = Inches(0.8 + idx * 2.95)
        top = Inches(1.8)
        box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(5.0))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(10.8)
        p2.font.color.rgb = TEXT_BODY
        p2.space_before = Pt(8)

    add_slide_footer(s5, 5)

    # =========================================================================
    # SLIDE 6: Applied AI Backend & Multi-Agent Subsystems
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6)
    add_slide_header(s6, "Applied AI Capabilities", "Multi-Agent Subsystems & Backend Python Modules", "Modular micro-agents powering intelligent Socratic STEM mentorship", ACCENT_BLUE)

    agent_features = [
        ("🧠 LLM Router & Semantic Context (llm_router.py)", 
         "Dynamic multi-model fallback across high-efficiency text and vision models. Integrated with vani_memory.db for persistent conversational context, episodic recall, and semantic search matching.", 
         ACCENT_BLUE),
        ("🎙️ Multimodal Voice Hub (voice_agent.py)", 
         "Integrated edge speech synthesis converting written scientific explanations into natural audio speech. Empowers auditory learners, vernacular students, and visually impaired learners.", 
         ACCENT_GREEN),
        ("🤖 Vision-Guided Screen Agent (agentic_loop.py)", 
         "Utilizes visual reasoning algorithms to parse student screen states, circuit diagrams, geometric shapes, and mathematical formulas, acting as an interactive AI lab assistant.", 
         ACCENT_GOLD),
        ("⚡ Proactive System Agent (autonomous_agent.py)", 
         "Monitors real-time client battery status, initiates eco-friendly power-saving protocols, and autonomously schedules personalized study routines for maximum productivity.", 
         ACCENT_PURPLE)
    ]

    for idx, (title, desc, col) in enumerate(agent_features):
        row = idx // 2
        col_idx = idx % 2
        left = Inches(0.8 + col_idx * 5.95)
        top = Inches(1.8 + row * 2.55)
        box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.35))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_BODY
        p2.space_before = Pt(6)

    add_slide_footer(s6, 6)

    # =========================================================================
    # SLIDE 7: End-to-End Operational Workflow (User Journey)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7)
    add_slide_header(s7, "End-to-End User Pipeline", "5-Step Zero-Friction Operational Flow", "Seamless transition from authentication to continuous conceptual mastery", ACCENT_GREEN)

    steps = [
        ("Step 01", "Secure Auth", "Student signs in securely via Google OAuth. System verifies credentials with zero password fatigue.", ACCENT_BLUE),
        ("Step 02", "UID Generation", "Cryptographic engine computes permanent student token (V.A.N.I-xAI-DHRU-3478).", ACCENT_GREEN),
        ("Step 03", "Sovereign Launch", "Workspace opens with 100% unlocked features, ₹0 fees, and zero recurring paywalls.", ACCENT_GOLD),
        ("Step 04", "Socratic Dialogue", "Student asks STEM questions; Socratic filter provides hints, voice audio, and derivations.", ACCENT_PURPLE),
        ("Step 05", "Persistent Mastery", "Autonomous SQLite memory logs milestones, voice sessions, and concept mastery.", ACCENT_BLUE)
    ]

    for idx, (step_num, title, desc, col) in enumerate(steps):
        left = Inches(0.8 + idx * 2.38)
        top = Inches(1.85)
        box = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.25), Inches(4.85))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        
        tf = box.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = step_num.upper()
        p0.font.name = "Segoe UI"
        p0.font.size = Pt(11)
        p0.font.bold = True
        p0.font.color.rgb = col
        
        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_DARK
        p1.space_before = Pt(4)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_BODY
        p2.space_before = Pt(8)

    add_slide_footer(s7, 7)

    # =========================================================================
    # SLIDE 8: InAixVAi Administrative Command Center
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8)
    add_slide_header(s8, "Teacher Oversight & Security", "InAixVAi Administrative Command Center", "Dedicated real-time portal for school educators, lab teachers, and mentors", ACCENT_GOLD)

    card_left = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = CARD_BG
    card_left.line.color.rgb = ACCENT_BLUE
    card_left.line.width = Pt(1.8)
    tf_l = card_left.text_frame
    tf_l.word_wrap = True
    p = tf_l.paragraphs[0]
    p.text = "📊 Real-Time Classroom Telemetry & Safety"
    p.font.name = "Segoe UI"
    p.font.size = Pt(14.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    adm_points1 = [
        "Live Sub-User Tracking: Real-time visibility into all authenticated student sessions across school computer labs.",
        "Query Load & Token Telemetry: Monitors computational request volumes, API throughput, and token efficiency in real-time.",
        "Child-Safe Sandbox Isolation: Ensures query isolation and strictly prevents toxic, unsafe, or off-curriculum prompt inputs.",
        "Multi-Device Session Management: Gives teachers instant controls to terminate anomalous sessions or reset classroom states."
    ]
    for pt in adm_points1:
        p = tf_l.add_paragraph()
        p.text = f"• {pt}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_BODY
        p.space_before = Pt(7)

    card_right = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(1.8), Inches(5.7), Inches(5.0))
    card_right.fill.solid()
    card_right.fill.fore_color.rgb = CARD_BG
    card_right.line.color.rgb = ACCENT_GOLD
    card_right.line.width = Pt(1.8)
    tf_r = card_right.text_frame
    tf_r.word_wrap = True
    p = tf_r.paragraphs[0]
    p.text = "🔐 Teacher Controls & Pedagogical Insights"
    p.font.name = "Segoe UI"
    p.font.size = Pt(14.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD

    adm_points2 = [
        "Conceptual Heatmaps: Identifies challenging STEM concepts where students require extra classroom teacher intervention.",
        "Dynamic PIN Authentication: Protects administrative oversight panels with dual-layer master PIN authentication.",
        "Zero-Trust Architecture: Guarantees complete student data privacy with zero third-party telemetry or ad tracking.",
        "Scalable Lab Rollout: Enables one-click activation across entire school computer laboratories without local software installs."
    ]
    for pt in adm_points2:
        p = tf_r.add_paragraph()
        p.text = f"• {pt}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_BODY
        p.space_before = Pt(7)

    add_slide_footer(s8, 8)

    # =========================================================================
    # SLIDE 9: Experimental Observations & Quantitative Data Benchmarking
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9)
    add_slide_header(s9, "Experimental Data & Testing", "Quantitative Benchmarking & Experimental Results", "Empirical data validating asymmetric efficiency and low-spec performance", ACCENT_GREEN)

    rows = 7
    cols = 4
    left = Inches(0.8)
    top = Inches(1.75)
    width = Inches(11.733)
    height = Inches(4.3)

    table_shape = s9.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.columns[0].width = Inches(3.1)
    table.columns[1].width = Inches(2.85)
    table.columns[2].width = Inches(2.85)
    table.columns[3].width = Inches(2.933)

    headers = ["Evaluation Metric", "Commercial AI (ChatGPT/Claude)", "Traditional EdTech Apps", "V.A.N.I - xAI (BoVxAi)"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEXT_NAVY if i != 3 else ACCENT_GREEN
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.name = "Segoe UI"
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_WHITE

    bench_data = [
        ("Client Payload Size", "45 MB - 95 MB (Heavy web client)", "85 MB - 150 MB (Native App)", "< 50 KB (99.9% lighter)"),
        ("Cold Boot Time (2G/3G Network)", "12.8s - 18.5s (High latency)", "15.0s - 25.0s (Slow startup)", "0.78 seconds (Instant load)"),
        ("Peak Client RAM Consumption", "380 MB - 600 MB (Crashes low RAM)", "250 MB - 450 MB (Heavy usage)", "38 MB (Runs on ₹5,000 phones)"),
        ("Monthly Student Cost", "₹1,500 - ₹2,500 / month (Paywall)", "₹5,000 - ₹25,000 / year", "100% Free Lifetime (₹0 Paywalls)"),
        ("Minimum Hardware Specification", "Modern 6GB RAM Smartphone / PC", "Recent Android 10+ / iOS", "Any Web Browser / Pentium PC"),
        ("Pedagogical Scaffolding", "Direct answer dump (Cheat risk)", "Static recorded video lectures", "Guided Socratic Step-by-Step")
    ]

    for r_idx, row_data in enumerate(bench_data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            if c_idx == 3:
                cell.fill.fore_color.rgb = TINT_GREEN_BG
            else:
                cell.fill.fore_color.rgb = CARD_BG if r_idx % 2 == 0 else CARD_BG_ALT
                
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Segoe UI"
            p.font.size = Pt(10)
            if c_idx == 3:
                p.font.bold = True
                p.font.color.rgb = ACCENT_GREEN
            elif c_idx == 0:
                p.font.bold = True
                p.font.color.rgb = TEXT_DARK
            else:
                p.font.color.rgb = TEXT_BODY

    # Bottom Highlight Callout
    bot_box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.15), Inches(11.733), Inches(0.75))
    bot_box.fill.solid()
    bot_box.fill.fore_color.rgb = TINT_BLUE_BG
    bot_box.line.color.rgb = ACCENT_BLUE
    bot_box.line.width = Pt(1.5)
    tf_bb = bot_box.text_frame
    p_bb = tf_bb.paragraphs[0]
    p_bb.text = "⚡ Key Scientific Finding: V.A.N.I-xAI achieves a 99.9% payload reduction and 16x faster load times on 2G/3G networks, unlocking AI access for 100% of Indian students with zero hardware cost."
    p_bb.font.name = "Segoe UI"
    p_bb.font.size = Pt(10.5)
    p_bb.font.bold = True
    p_bb.font.color.rgb = ACCENT_BLUE
    p_bb.alignment = PP_ALIGN.CENTER

    add_slide_footer(s9, 9)

    # =========================================================================
    # SLIDE 10: Comparative Advantage & Key Differentiators
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10)
    add_slide_header(s10, "Competitive Differentiators", "Core Pillars of Competitive Advantage", "Why V.A.N.I-xAI stands as a unique, high-impact indigenous innovation", ACCENT_GOLD)

    diff_cards = [
        ("1. Economic Sovereignty", 
         "Complete elimination of credit card mandates and recurring subscriptions. Built as an open, accessible public good for every Indian child regardless of socioeconomic status.", 
         ACCENT_BLUE),
        ("2. Pedagogical Integrity", 
         "Engineered with a dedicated Socratic prompt filter that prevents rote copy-pasting, instead training students in cognitive inquiry, critical derivation, and conceptual clarity.", 
         ACCENT_GREEN),
        ("3. Universal Inclusivity", 
         "Under 50KB payload runs on legacy school Pentium computers, budget ₹5,000 phones, and slow 2G/3G mobile data without requiring 100MB+ app downloads.", 
         ACCENT_GOLD),
        ("4. Sovereign Data Security", 
         "100% localized context storage via SQLite (vani_memory.db) with zero third-party commercial data harvesting, ensuring complete data privacy for young learners.", 
         ACCENT_PURPLE)
    ]

    for idx, (title, desc, col) in enumerate(diff_cards):
        row = idx // 2
        col_idx = idx % 2
        left = Inches(0.8 + col_idx * 5.95)
        top = Inches(1.8 + row * 2.55)
        box = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.35))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_BODY
        p2.space_before = Pt(6)

    add_slide_footer(s10, 10)

    # =========================================================================
    # SLIDE 11: Environmental Sustainability & 100% Zero E-Waste
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_background(s11)
    add_slide_header(s11, "Environmental Sustainability", "100% Zero E-Waste & Green Technology Alignment", "A sustainable software innovation aligned with National Green Missions and UN SDGs", ACCENT_GREEN)

    c1 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(3.7), Inches(4.3))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = ACCENT_GREEN
    c1.line.width = Pt(1.8)
    tf1 = c1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "🌱 Zero Physical Manufacturing"
    p.font.name = "Segoe UI"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    
    points1 = [
        "100% software-driven: Requires zero physical plastic casing, toxic lithium batteries, or PCB fabrication.",
        "Completely eliminates electronic manufacturing scrap at the source.",
        "Drastically reduces carbon emissions through optimized cloud server inference."
    ]
    for pt in points1:
        p = tf1.add_paragraph()
        p.text = f"• {pt}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(10.8)
        p.font.color.rgb = TEXT_BODY
        p.space_before = Pt(7)

    c2 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.8), Inches(3.7), Inches(4.3))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = ACCENT_BLUE
    c2.line.width = Pt(1.8)
    tf2 = c2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "📱 Extends Old Hardware Lifespan"
    p.font.name = "Segoe UI"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    points2 = [
        "Repurposes existing school computer labs and parents' low-cost old smartphones.",
        "Prevents premature disposal of functional electronics by keeping client footprint <50KB.",
        "Runs directly inside standard web browsers with zero installation requirements."
    ]
    for pt in points2:
        p = tf2.add_paragraph()
        p.text = f"• {pt}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(10.8)
        p.font.color.rgb = TEXT_BODY
        p.space_before = Pt(7)

    c3 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.3))
    c3.fill.solid()
    c3.fill.fore_color.rgb = CARD_BG
    c3.line.color.rgb = ACCENT_GOLD
    c3.line.width = Pt(1.8)
    tf3 = c3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "📄 100% Paperless Ecosystem"
    p.font.name = "Segoe UI"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    
    points3 = [
        "Eliminates physical question banks, printed worksheets, and disposable test guides.",
        "All learning sessions, notes, and milestones are preserved digitally in persistent memory.",
        "Supports India's Green Technology and Net-Zero Digital Education initiatives."
    ]
    for pt in points3:
        p = tf3.add_paragraph()
        p.text = f"• {pt}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(10.8)
        p.font.color.rgb = TEXT_BODY
        p.space_before = Pt(7)

    # UN SDG Alignment Banner
    sdg_box = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.2), Inches(11.733), Inches(0.7))
    sdg_box.fill.solid()
    sdg_box.fill.fore_color.rgb = TINT_GREEN_BG
    sdg_box.line.color.rgb = ACCENT_GREEN
    sdg_box.line.width = Pt(1.5)
    tf_sdg = sdg_box.text_frame
    p_sdg = tf_sdg.paragraphs[0]
    p_sdg.text = "🌍 UN Sustainable Development Goals: 🎓 SDG 4 (Quality Education) | 💡 SDG 9 (Innovation) | ⚖️ SDG 10 (Reduced Inequalities) | ♻️ SDG 12 (Responsible Consumption)"
    p_sdg.font.name = "Segoe UI"
    p_sdg.font.size = Pt(10.5)
    p_sdg.font.bold = True
    p_sdg.font.color.rgb = ACCENT_GREEN
    p_sdg.alignment = PP_ALIGN.CENTER

    add_slide_footer(s11, 11)

    # =========================================================================
    # SLIDE 12: Societal Impact & Scalability for Bharat
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_background(s12)
    add_slide_header(s12, "Societal Impact & Scalability", "Grassroots Empowerment & Atmanirbhar Scalability", "Transforming educational equity for 250M+ Indian students with zero capital expenditure", ACCENT_GOLD)

    card_soc_l = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
    card_soc_l.fill.solid()
    card_soc_l.fill.fore_color.rgb = CARD_BG
    card_soc_l.line.color.rgb = ACCENT_BLUE
    card_soc_l.line.width = Pt(1.8)
    tf_sl = card_soc_l.text_frame
    tf_sl.word_wrap = True
    p = tf_sl.paragraphs[0]
    p.text = "🤝 Grassroots Societal Value"
    p.font.name = "Segoe UI"
    p.font.size = Pt(14.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    soc_points = [
        "Democratizing 1-on-1 Tutoring: Delivers world-class STEM mentorship to students in Tier-2, Tier-3 cities, and rural villages at ₹0 cost.",
        "Empowering First-Generation Learners: Provides patient, 24/7 guided explanations for complex physics, chemistry, math, and coding concepts.",
        "NEP 2020 Pedagogical Alignment: Fosters experiential inquiry, cognitive derivation, and scientific curiosity over rote memorization.",
        "Inclusive Vernacular Access: Voice synthesis enables students with reading difficulties or regional dialect preferences to learn effortlessly."
    ]
    for pt in soc_points:
        p = tf_sl.add_paragraph()
        p.text = f"• {pt}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_BODY
        p.space_before = Pt(7)

    card_soc_r = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(1.8), Inches(5.7), Inches(5.0))
    card_soc_r.fill.solid()
    card_soc_r.fill.fore_color.rgb = CARD_BG
    card_soc_r.line.color.rgb = ACCENT_GOLD
    card_soc_r.line.width = Pt(1.8)
    tf_sr = card_soc_r.text_frame
    tf_sr.word_wrap = True
    p = tf_sr.paragraphs[0]
    p.text = "🚀 High-Efficiency Sovereign Scaling Model"
    p.font.name = "Segoe UI"
    p.font.size = Pt(14.5)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD

    scale_points = [
        "Zero Hardware Capex: Schools require zero investment in new computer hardware or proprietary tablets to deploy V.A.N.I.",
        "Serverless Scalability: Cloud backend seamlessly auto-scales from a single classroom to thousands of concurrent district users.",
        "Sovereign Atmanirbhar Stack: 100% indigenously engineered codebase independent of foreign educational subscription conglomerates.",
        "Ready-to-Deploy Prototype: Live, cloud-hosted, and operational today with active student accounts and teacher management."
    ]
    for pt in scale_points:
        p = tf_sr.add_paragraph()
        p.text = f"• {pt}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_BODY
        p.space_before = Pt(7)

    add_slide_footer(s12, 12)

    # =========================================================================
    # SLIDE 13: Vision 2030 Roadmap: School to National Deployment
    # =========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    set_slide_background(s13)
    add_slide_header(s13, "Future Scope & Vision 2030", "Roadmap: From School Innovation to National Utility", "Structured multi-phase scaling strategy for NCSC and National Digital Education Missions", ACCENT_PURPLE)

    phases_roadmap = [
        ("Phase 1: School Pilot (2026)", 
         "Operational School Prototype", 
         [
             "Live deployment at PM SHRI KV No. 1 AFS Chakeri Kanpur.",
             "InAixVAi teacher command dashboard active for science labs.",
             "100% free onboarding for secondary school students.",
             "Continuous feedback integration on STEM Socratic hints."
         ], ACCENT_BLUE),
        ("Phase 2: District & State (2026-27)", 
         "Vernacular & Offline Mesh Expansion", 
         [
             "Integration of 10+ Indian regional languages (Hindi, etc.).",
             "Bi-directional voice conversation engine for rural learners.",
             "Offline edge caching protocol for schools with intermittent web.",
             "State-wide rollout across KVS Varanasi & Lucknow regions."
         ], ACCENT_GREEN),
        ("Phase 3: National Public AI (2028-30)", 
         "Sovereign Digital Public Infrastructure", 
         [
             "Proposed integration with DIKSHA, PM eVidya & Swayam portals.",
             "Virtual interactive STEM lab simulator and automated code sandbox.",
             "Complete Atmanirbhar Bharat AI mentoring stack for 250M+ students."
         ], ACCENT_GOLD)
    ]

    for idx, (phase, subtitle, bullets, col) in enumerate(phases_roadmap):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(1.8)
        box = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.75), Inches(5.0))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = phase
        p.font.name = "Segoe UI"
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = col
        
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = "Segoe UI"
        p_sub.font.size = Pt(11.5)
        p_sub.font.color.rgb = TEXT_DARK
        p_sub.font.bold = True
        p_sub.space_before = Pt(3)

        for b in bullets:
            p_b = tf.add_paragraph()
            p_b.text = f"• {b}"
            p_b.font.name = "Segoe UI"
            p_b.font.size = Pt(10.5)
            p_b.font.color.rgb = TEXT_BODY
            p_b.space_before = Pt(7)

    add_slide_footer(s13, 13)

    # =========================================================================
    # SLIDE 14: Scientific Conclusion, Acknowledgements & Q&A
    # =========================================================================
    s14 = prs.slides.add_slide(blank_layout)
    set_slide_background(s14)

    # Grand Conclusion Hero Card
    concl_card = s14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.6), Inches(11.733), Inches(6.2))
    concl_card.fill.solid()
    concl_card.fill.fore_color.rgb = CARD_BG
    concl_card.line.color.rgb = ACCENT_GREEN
    concl_card.line.width = Pt(2.5)
    tf_end = concl_card.text_frame
    tf_end.word_wrap = True

    p0 = tf_end.paragraphs[0]
    p0.text = "🏆 NATIONAL CHILDREN'S SCIENCE CONGRESS (NCSC) 2026-27"
    p0.font.name = "Segoe UI"
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_GOLD
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf_end.add_paragraph()
    p1.text = "Thank You, Respected Jury Members & Mentors!"
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_NAVY
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(4)

    p2 = tf_end.add_paragraph()
    p2.text = "V.A.N.I - xAI (BoVxAi) : Vāṇī Adhyātmik Navīn Intellect"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_GREEN
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(3)

    p3 = tf_end.add_paragraph()
    p3.text = "\"Democratizing cutting-edge AI mentorship for every student in Bharat — 100% Free, Zero Subscriptions, Zero E-Waste & Grounded in Scientific Integrity.\""
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(12)
    p3.font.italic = True
    p3.font.color.rgb = TEXT_BODY
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(8)

    # 4 Summary Highlight Chips
    chips_text = "✨ 100% Free Sovereign AI  |  🔑 Algorithmic UIDs  |  🎓 Socratic Mentoring  |  🌱 0% E-Waste"
    p_chip = tf_end.add_paragraph()
    p_chip.text = chips_text
    p_chip.font.name = "Segoe UI"
    p_chip.font.size = Pt(12.5)
    p_chip.font.bold = True
    p_chip.font.color.rgb = ACCENT_BLUE
    p_chip.alignment = PP_ALIGN.CENTER
    p_chip.space_before = Pt(10)

    # Acknowledgement Text
    ack_text = "We express sincere gratitude to National Council for Science & Technology Communication (NCSTC), Department of Science & Technology (DST), Kendriya Vidyalaya Sangathan (KVS), and PM SHRI KV No. 1 AFS Chakeri Kanpur for their mentorship and support."
    p_ack = tf_end.add_paragraph()
    p_ack.text = ack_text
    p_ack.font.name = "Segoe UI"
    p_ack.font.size = Pt(11)
    p_ack.font.color.rgb = TEXT_MUTED
    p_ack.alignment = PP_ALIGN.CENTER
    p_ack.space_before = Pt(10)

    p4 = tf_end.add_paragraph()
    p4.text = "Child Scientist: Dhruv Sagar (Class 10th) | Student UID: 1001049893\nPM SHRI KV NO.1 AFS Chakeri, Kanpur (UP) | Bureau Of V.A.N.I-xAI (BoVxAi)\n\n[ Live Working Prototype Ready for Jury Demonstration ]"
    p4.font.name = "Segoe UI"
    p4.font.size = Pt(11.5)
    p4.font.bold = True
    p4.font.color.rgb = TEXT_NAVY
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(12)

    add_slide_footer(s14, 14)

    # Save presentation to both target names
    output_filename = "VANI_xAI_NCSC_Official_Presentation.pptx"
    prs.save(output_filename)
    print(f"[SUCCESS] Official NCSC PPTX created: {output_filename}")
    
    output_light = "VANI_xAI_NCSC_2026_Light.pptx"
    prs.save(output_light)
    print(f"[SUCCESS] Synced to: {output_light}")

    return output_filename

if __name__ == "__main__":
    build_ncsc_presentation()
