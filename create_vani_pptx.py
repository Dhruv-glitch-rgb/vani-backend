import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_vani_presentation():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Theme Colors - Aligned with VANI Project Report PDF UI
    BG_CANVAS = RGBColor(248, 250, 252)     # #F8FAFC (Clean Light Slate)
    CARD_BG = RGBColor(255, 255, 255)       # #FFFFFF (Crisp White Card)
    CARD_BG_ALT = RGBColor(241, 245, 249)   # #F1F5F9 (Soft Slate)
    CARD_BORDER = RGBColor(226, 232, 240)   # #E2E8F0 (Subtle Slate Border)
    HERO_DARK = RGBColor(11, 15, 25)        # #0B0F19 (Deep Obsidian Header)
    HERO_CARD_BG = RGBColor(26, 36, 56)     # #1A2438 (Dark Slate Card)
    
    TEXT_DARK = RGBColor(15, 23, 42)        # #0F172A (Deep Slate Navy)
    TEXT_BODY = RGBColor(51, 65, 85)        # #334155 (Slate Charcoal)
    TEXT_MUTED = RGBColor(100, 116, 139)    # #64748B (Slate Muted)
    TEXT_WHITE = RGBColor(248, 250, 252)    # #F8FAFC (Pure White)
    
    ACCENT_BLUE = RGBColor(2, 132, 199)     # #0284C7 (Sky/Primary Blue)
    ACCENT_CYAN = RGBColor(56, 189, 248)    # #38BDF8 (Electric Cyan)
    ACCENT_GREEN = RGBColor(22, 163, 74)    # #16A34A (Emerald Green)
    ACCENT_MINT = RGBColor(52, 211, 153)    # #34D399 (Mint Green)
    ACCENT_GOLD = RGBColor(217, 119, 6)     # #D97706 (Amber/Gold)
    ACCENT_YELLOW = RGBColor(251, 191, 36)  # #FBBF24 (Bright Gold)
    ACCENT_PURPLE = RGBColor(147, 51, 234)  # #9333EA (Purple/Violet)
    ACCENT_RED = RGBColor(220, 38, 38)      # #DC2626 (Coral Red)

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_CANVAS
        bg.line.fill.background()
        return bg

    def add_header(slide, category_tag, title_text, subtitle_text=None, tag_color=ACCENT_BLUE):
        # Category Tag badge
        tag_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.42), Inches(3.8), Inches(0.35))
        tag_box.fill.solid()
        tag_box.fill.fore_color.rgb = CARD_BG
        tag_box.line.color.rgb = tag_color
        tag_box.line.width = Pt(1.5)
        tf = tag_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_tag.upper()
        p.font.name = "Segoe UI"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = tag_color
        p.alignment = PP_ALIGN.CENTER

        # Title Textbox
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.82), Inches(11.7), Inches(0.8))
        tf_t = tb.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.name = "Segoe UI"
        p_t.font.size = Pt(25)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_DARK

        if subtitle_text:
            p_s = tf_t.add_paragraph()
            p_s.text = subtitle_text
            p_s.font.name = "Segoe UI"
            p_s.font.size = Pt(12)
            p_s.font.color.rgb = TEXT_MUTED
            p_s.space_before = Pt(3)

    def add_footer(slide, slide_num):
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.35))
        tf_f = footer_box.text_frame
        p_f = tf_f.paragraphs[0]
        p_f.text = f"V.A.N.I - xAI (BoVxAi) | National Children's Science Congress (NCSC) | Slide {slide_num:02d}"
        p_f.font.name = "Segoe UI"
        p_f.font.size = Pt(10)
        p_f.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 1: Title Slide (Grand Executive Showcase)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s1)

    # Master Dark Hero Card
    hero_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.5), Inches(11.733), Inches(3.2))
    hero_card.fill.solid()
    hero_card.fill.fore_color.rgb = HERO_DARK
    hero_card.line.color.rgb = ACCENT_CYAN
    hero_card.line.width = Pt(2)
    tf_hero = hero_card.text_frame
    tf_hero.word_wrap = True

    # Top National Badge inside Hero
    p0 = tf_hero.paragraphs[0]
    p0.text = "🏆 NATIONAL CHILDREN'S SCIENCE CONGRESS (NCSC) 2026-27"
    p0.font.name = "Segoe UI"
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_YELLOW
    p0.alignment = PP_ALIGN.LEFT

    p1 = tf_hero.add_paragraph()
    p1.text = "V.A.N.I - xAI (BoVxAi)"
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(38)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_CYAN
    p1.alignment = PP_ALIGN.LEFT
    p1.space_before = Pt(2)

    p2 = tf_hero.add_paragraph()
    p2.text = "Vāṇī Adhyātmik Navīn Intellect & Bridge of Voice"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(17)
    p2.font.color.rgb = ACCENT_MINT
    p2.font.bold = True
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(2)

    p3 = tf_hero.add_paragraph()
    p3.text = "Next-Generation 100% Free Sovereign AI Tutoring, Multimodal Voice & Autonomous Learning Ecosystem for Bharat"
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(12.5)
    p3.font.color.rgb = TEXT_WHITE
    p3.alignment = PP_ALIGN.LEFT
    p3.space_before = Pt(4)

    # Embed Official Bureau of VANI-xAI Logo on Right Side of Hero Card
    official_logo = "svg.png"
    if not os.path.exists(official_logo):
        official_logo = os.path.join("VANI-B_IMGs", "VANI_xAI_Official_Logo.png")

    if os.path.exists(official_logo):
        box1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(0.7), Inches(2.5), Inches(2.75))
        box1.fill.solid()
        box1.fill.fore_color.rgb = CARD_BG
        box1.line.color.rgb = ACCENT_CYAN
        box1.line.width = Pt(1.5)
        s1.shapes.add_picture(official_logo, Inches(10.0), Inches(0.85), width=Inches(2.1))

    # Info Cards Container (White cards below Hero)
    card1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.9), Inches(5.7), Inches(2.95))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = ACCENT_BLUE
    card1.line.width = Pt(1.8)
    tf_c1 = card1.text_frame
    tf_c1.word_wrap = True
    c1_p1 = tf_c1.paragraphs[0]
    c1_p1.text = "👤 Innovator Profile"
    c1_p1.font.name = "Segoe UI"
    c1_p1.font.size = Pt(16)
    c1_p1.font.bold = True
    c1_p1.font.color.rgb = ACCENT_BLUE

    items_c1 = [
        ("Student Name:", "Dhruv Sagar"),
        ("Class & Stream:", "Class 10th (Secondary Section)"),
        ("Institution:", "PM SHRI KV NO.1 AFS Chakeri, Kanpur"),
        ("State / Region:", "Uttar Pradesh, India (KVS RO Varanasi)"),
        ("Department:", "Dept. of Science & Technology (DST) / NIF India")
    ]
    for k, v in items_c1:
        p = tf_c1.add_paragraph()
        p.text = f"• {k} {v}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_BODY
        p.space_before = Pt(4)

    card2 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.833), Inches(3.9), Inches(5.7), Inches(2.95))
    card2.fill.solid()
    card2.fill.fore_color.rgb = CARD_BG
    card2.line.color.rgb = ACCENT_GREEN
    card2.line.width = Pt(1.8)
    tf_c2 = card2.text_frame
    tf_c2.word_wrap = True
    c2_p1 = tf_c2.paragraphs[0]
    c2_p1.text = "🎯 Project Key Attributes"
    c2_p1.font.name = "Segoe UI"
    c2_p1.font.size = Pt(16)
    c2_p1.font.bold = True
    c2_p1.font.color.rgb = ACCENT_GREEN

    items_c2 = [
        ("Innovation Domain:", "Information Technology / Applied Artificial Intelligence"),
        ("Core Model:", "100% Free Sovereign AI (Zero Subscriptions, Zero Paywalls)"),
        ("Architecture:", "Asymmetric Cloud-Native Sub-User Engine (<50KB Client)"),
        ("Sustainability:", "100% Zero E-Waste & Zero Plastic Hardware"),
        ("Prototype Status:", "Fully Operational & Deployed Live on Web")
    ]
    for k, v in items_c2:
        p = tf_c2.add_paragraph()
        p.text = f"• {k} {v}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_BODY
        p.space_before = Pt(4)

    add_footer(s1, 1)

    # =========================================================================
    # SLIDE 2: Grassroots Problem Statement
    # =========================================================================
    s2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s2)
    add_header(s2, "Grassroots Problem Identification", "The Educational & Technological Divide in India", "Why generic commercial AI models fail 90% of Indian school students", ACCENT_RED)

    problems = [
        ("💸 Severe Economic Exclusivity & Paywalls", 
         "Commercial AI tools (ChatGPT Plus, Claude Pro) charge ₹1,500 to ₹2,500/month and mandate recurring subscriptions. This locks out millions of talented underprivileged Indian students behind paywalls.", 
         ACCENT_GOLD),
        ("📱 Hardware & Low-Bandwidth Barrier", 
         "Advanced generative models demand high-end GPUs, latest smartphones, and high-speed broadband. Underprivileged students using ₹5,000 smartphones and 2G/3G mobile data cannot run heavy AI apps.", 
         ACCENT_BLUE),
        ("🧠 The 'Cheat-Engine' Flaw", 
         "Mainstream LLMs generate direct solutions without explanations. This encourages rote copy-pasting and academic cheating rather than cognitive reasoning, critical inquiry, or conceptual mastery.", 
         ACCENT_RED),
        ("🛡️ Lack of Localized & Safe Guardrails", 
         "Foreign AI platforms lack alignment with Indian NCERT/CBSE curriculum, Vernacular voice capabilities, and youth-safe ethical guardrails tailored for Indian schools and children.", 
         ACCENT_PURPLE)
    ]

    for idx, (head, desc, col) in enumerate(problems):
        row = idx // 2
        col_idx = idx % 2
        left = Inches(0.8 + col_idx * 5.95)
        top = Inches(1.8 + row * 2.55)
        box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.35))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = head
        p.font.name = "Segoe UI"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_BODY
        p2.space_before = Pt(8)

    add_footer(s2, 2)

    # =========================================================================
    # SLIDE 3: The Innovation - V.A.N.I-xAI & BoVxAi Engine
    # =========================================================================
    s3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s3)
    add_header(s3, "The Proposed Innovation", "V.A.N.I-xAI: Sovereign AI Mentor for Every Student", "Democratizing cloud intelligence through indigenous asymmetric architecture", ACCENT_GREEN)

    solutions = [
        ("⚡ Asymmetric Sub-User Core", 
         "99% of computation is shifted to high-speed cloud clusters. The web client payload is <50KB, allowing instant sub-second boot times even on legacy ₹5,000 mobile phones and slow 2G/3G networks.",
         ACCENT_BLUE),
        ("🎓 Socratic Mentorship Engine", 
         "Instead of regurgitating answers, the model acts as an empathetic digital mentor—breaking complex STEM concepts into interactive hints, guided questioning, and conceptual scaffolding.",
         ACCENT_GREEN),
        ("🎙️ BoV Multimodal Voice Hub", 
         "Integrated speech synthesis and voice processing (`voice_agent.py`) enables real-time auditory explanations, assisting students with reading difficulties or vernacular preferences.",
         ACCENT_PURPLE),
        ("🔓 100% Free Sovereign Access (Zero Paywalls)", 
         "Removed all premium subscription plans and paywalls. Every capability is 100% free and permanently unlocked for all students with seamless Google OAuth and mathematical UID verification.",
         ACCENT_GOLD)
    ]

    for idx, (title, desc, col) in enumerate(solutions):
        row = idx // 2
        col_idx = idx % 2
        left = Inches(0.8 + col_idx * 5.95)
        top = Inches(1.8 + row * 2.55)
        box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.35))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"0{idx+1}. {title}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_BODY
        p2.space_before = Pt(8)

    add_footer(s3, 3)

    # =========================================================================
    # SLIDE 4: Scientific Novelty & 4 Core Pillars
    # =========================================================================
    s4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s4)
    add_header(s4, "Scientific Novelty & IP", "4 Pillars of Novel Technological Architecture", "Handcrafted algorithmic systems engineered from scratch for the V.A.N.I ecosystem", ACCENT_GOLD)

    pillars = [
        ("1. Mathematical UID Generation", 
         "Algorithmically generates permanent user-bound tokens (e.g., `V.A.N.I-xAI-DHRU-3478`) upon initial Google OAuth sign-in. This permanently locks memory and workspace states to verified student identities without hardware lock-in.", 
         ACCENT_BLUE),
        ("2. 3-Layer Cryptographic Security", 
         "Proprietary algorithmic authorization validated against 3 rigorous integrity checks: (1) Database Existence, (2) Active Student Token validation, and (3) Strict UID binding to eliminate impersonation.", 
         ACCENT_GREEN),
        ("3. InAixVAi Command Dashboard", 
         "A bespoke educator/admin dashboard enabling real-time multi-user telemetry, classroom activity monitoring, query load oversight, and platform integrity management with zero subscription barriers.", 
         ACCENT_GOLD),
        ("4. Autonomous Session & Memory Protocol", 
         "An autonomous background self-healing engine that manages real-time cross-device session synchronization, SQLite memory persistence (`vani_memory.db`), and load-balancing with zero maintenance.", 
         ACCENT_PURPLE)
    ]

    for idx, (title, desc, col) in enumerate(pillars):
        left = Inches(0.8 + idx * 2.95)
        top = Inches(1.8)
        box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(5.0))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = TEXT_BODY
        p2.space_before = Pt(10)

    add_footer(s4, 4)

    # =========================================================================
    # SLIDE 5: BoVxAi Autonomous Intelligence & Multi-Agent Loop
    # =========================================================================
    s5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s5)
    add_header(s5, "Advanced Applied AI Capabilities", "BoVxAi: Multimodal Voice, Memory & Agentic Control", "State-of-the-art intelligent modules integrated into the local & cloud backend", ACCENT_BLUE)

    agent_features = [
        ("🧠 LLM Router & SQLite Memory", 
         "Features `llm_router.py` with multi-model fallback across vision and text models. Backed by `vani_memory.db` for persistent conversational memory, contextual recall, and fast semantic search.", 
         ACCENT_BLUE),
        ("🎙️ Multimodal Voice Engine", 
         "Powered by `voice_agent.py` and edge TTS. Converts complex scientific derivations and explanations into natural audio speech for inclusive hands-free learning.", 
         ACCENT_GREEN),
        ("🤖 Vision-Guided Agentic Loop", 
         "`agentic_loop.py` utilizes visual reasoning to analyze user screen states, coordinates, and automated workflows, transforming V.A.N.I into an interactive digital lab assistant.", 
         ACCENT_GOLD),
        ("⚡ Proactive System Agent", 
         "`autonomous_agent.py` monitors real-time battery status, initiates power-saving protocols, and schedules automated learning routines for student productivity.", 
         ACCENT_PURPLE)
    ]

    for idx, (title, desc, col) in enumerate(agent_features):
        row = idx // 2
        col_idx = idx % 2
        left = Inches(0.8 + col_idx * 5.95)
        top = Inches(1.8 + row * 2.55)
        box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.35))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_BODY
        p2.space_before = Pt(8)

    add_footer(s5, 5)

    # =========================================================================
    # SLIDE 6: End-to-End System Architecture & Data Pipeline
    # =========================================================================
    s6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s6)
    add_header(s6, "End-to-End Pipeline", "5-Step Operational Architecture & User Workflow", "Zero-friction, 100% free sovereign workflow from authentication to Socratic response", ACCENT_GREEN)

    steps = [
        ("Step 1", "Auth & UID Gen", "Student signs in securely via Google OAuth. System algorithmically computes and permanently assigns their unique V.A.N.I UID.", ACCENT_BLUE),
        ("Step 2", "Instant Free Access", "Student enters workspace with 100% unlocked features, zero subscription barriers, zero fees, and instant workspace launch.", ACCENT_GREEN),
        ("Step 3", "Security & Swarm Sync", "System cryptographically binds UID with secure session tokens, enabling instant cross-device pairing and encrypted memory persistence.", ACCENT_GOLD),
        ("Step 4", "Socratic AI Inference", "Queries route through the LLM router, memory database, and Socratic safety filters for guided step-by-step tutoring.", ACCENT_PURPLE),
        ("Step 5", "Continuous Learning", "Autonomous memory engine continuously stores learning milestones, voice history, and concept mastery for lifelong growth.", ACCENT_BLUE)
    ]

    for idx, (step_num, title, desc, col) in enumerate(steps):
        left = Inches(0.8 + idx * 2.38)
        top = Inches(1.9)
        box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.25), Inches(4.8))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        
        tf = box.text_frame
        tf.word_wrap = True
        
        # Step number badge
        p0 = tf.paragraphs[0]
        p0.text = step_num.upper()
        p0.font.name = "Segoe UI"
        p0.font.size = Pt(12)
        p0.font.bold = True
        p0.font.color.rgb = col
        
        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(13.5)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_DARK
        p1.space_before = Pt(6)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_BODY
        p2.space_before = Pt(10)

    add_footer(s6, 6)

    # =========================================================================
    # SLIDE 7: Comparative Advantage Matrix
    # =========================================================================
    s7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s7)
    add_header(s7, "Competitive Benchmarking", "V.A.N.I-xAI vs. Conventional AI & EdTech Portals", "Engineered specifically for maximum grassroots accessibility, affordability & pedagogical integrity", ACCENT_GOLD)

    rows = 6
    cols = 4
    left = Inches(0.8)
    top = Inches(1.8)
    width = Inches(11.733)
    height = Inches(4.9)

    table_shape = s7.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.columns[0].width = Inches(3.1)
    table.columns[1].width = Inches(2.85)
    table.columns[2].width = Inches(2.85)
    table.columns[3].width = Inches(2.933)

    headers = ["Evaluation Metric", "Commercial AI (ChatGPT/Claude)", "Traditional EdTech Apps", "V.A.N.I - xAI (Our Innovation)"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = HERO_DARK if i != 3 else ACCENT_BLUE
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.name = "Segoe UI"
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_WHITE

    data = [
        ("Cost for Students", "₹1,500 - ₹2,500/month (Expensive)", "₹5,000 - ₹25,000/year (High cost)", "100% Free Lifetime (Zero Subscriptions)"),
        ("Hardware / Device Barrier", "Requires modern smartphone & RAM", "Requires 100MB+ app download", "Runs on any browser / <50KB payload"),
        ("Pedagogical Approach", "Direct answer generation (cheat)", "Static recorded video lectures", "Socratic guidance & Concept scaffolding"),
        ("Payment & Access Barrier", "Mandates International Credit Cards", "Closed recurring subscriptions", "Zero Paywalls / Instant Free Access"),
        ("Zero E-Waste & Ecology", "Heavy local battery & GPU drain", "Proprietary locked tablet hardware", "100% Cloud-Native (Zero E-Waste)")
    ]

    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            if c_idx == 3:
                cell.fill.fore_color.rgb = RGBColor(240, 253, 244) # Soft Emerald Tint
            else:
                cell.fill.fore_color.rgb = CARD_BG if r_idx % 2 == 0 else CARD_BG_ALT
                
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Segoe UI"
            p.font.size = Pt(11)
            if c_idx == 3:
                p.font.bold = True
                p.font.color.rgb = ACCENT_GREEN
            elif c_idx == 0:
                p.font.bold = True
                p.font.color.rgb = TEXT_DARK
            else:
                p.font.color.rgb = TEXT_BODY

    add_footer(s7, 7)

    # =========================================================================
    # SLIDE 8: Social Impact, Sustainability & DBT ₹10,000 Utilization
    # =========================================================================
    s8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s8)
    add_header(s8, "Societal Value & Scalability", "Grassroots Impact, Ecology & Resource Plan", "Maximizing grassroots educational impact for Bharat", ACCENT_GREEN)

    # 3 Large Cards
    c1 = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5.0))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = ACCENT_BLUE
    c1.line.width = Pt(1.8)
    tf1 = c1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "🌱 100% Zero E-Waste"
    p.font.name = "Segoe UI"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    points1 = [
        "100% software-driven: Requires zero physical plastic casing, toxic lithium batteries, or PCB manufacturing.",
        "Repurposes existing school computer labs and parents' low-cost smartphones.",
        "Drastically reduces carbon emissions by optimizing lightweight server inference over heavy local processing.",
        "Eco-friendly, completely paperless, and sustainable digital architecture."
    ]
    for pt in points1:
        p = tf1.add_paragraph()
        p.text = f"• {pt}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_BODY
        p.space_before = Pt(8)

    c2 = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5.0))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = ACCENT_GREEN
    c2.line.width = Pt(1.8)
    tf2 = c2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "🤝 Grassroots Empowerment"
    p.font.name = "Segoe UI"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    
    points2 = [
        "Delivers 24/7 personalized 1-on-1 STEM mentoring to Tier-2, Tier-3 & rural learners with zero subscription fees.",
        "Empowers first-generation learners with conceptual clarity in Science, Math, and Coding.",
        "Directly advances National Education Policy (NEP 2020) goals for experiential and equitable learning.",
        "Bridges the digital divide between metropolitan and village schools."
    ]
    for pt in points2:
        p = tf2.add_paragraph()
        p.text = f"• {pt}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_BODY
        p.space_before = Pt(8)

    c3 = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.0))
    c3.fill.solid()
    c3.fill.fore_color.rgb = CARD_BG
    c3.line.color.rgb = ACCENT_GOLD
    c3.line.width = Pt(1.8)
    tf3 = c3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "💰 ₹10,000 DBT Fund Budget"
    p.font.name = "Segoe UI"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    
    points3 = [
        "₹5,000: High-speed Cloud API credits (empowers 100,000+ student educational queries at zero cost to students).",
        "₹2,500: Serverless database, hosting & low-latency Indian Edge CDN scaling.",
        "₹1,500: Vernacular Language token optimization for Hindi & regional Indian dialects.",
        "₹1,000: SSL security, domain infrastructure & school cluster deployment."
    ]
    for pt in points3:
        p = tf3.add_paragraph()
        p.text = f"• {pt}"
        p.font.name = "Segoe UI"
        p.font.size = Pt(11.5)
        p.font.color.rgb = TEXT_BODY
        p.space_before = Pt(8)

    add_footer(s8, 8)

    # =========================================================================
    # SLIDE 9: Live Prototype Readiness & Validation
    # =========================================================================
    s9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s9)
    add_header(s9, "Working Prototype Demonstration", "Fully Operational & Deployed Features", "Tested and verified codebase ready for school & district implementation", ACCENT_BLUE)

    features = [
        ("🔐 InAixVAi Command Dashboard", 
         "A live administrative panel enabling teachers and admins to track authenticated sub-users, monitor query loads, manage classroom sessions, and safeguard student data in real-time.", 
         ACCENT_BLUE),
        ("🧠 LLM Router & Semantic Context", 
         "Equipped with multi-model fallback, local memory database (`vani_memory.db`), and Socratic prompt pipelines to ensure accurate, hallucination-resistant educational mentorship.", 
         ACCENT_GREEN),
        ("🎙️ Multimodal Voice Tutor (BoV)", 
         "Integrated voice synthesizer (`voice_agent.py`) converting written solutions to spoken words, enabling auditory learners and visually impaired students to engage effortlessly.", 
         ACCENT_GOLD),
        ("📱 Universal Cross-Device Compatibility", 
         "Zero-install, browser-accessible interface tested across low-end Android phones, iPads, Windows PCs, and Linux school lab workstations with <1s load latency.", 
         ACCENT_PURPLE)
    ]

    for idx, (title, desc, col) in enumerate(features):
        row = idx // 2
        col_idx = idx % 2
        left = Inches(0.8 + col_idx * 5.95)
        top = Inches(1.8 + row * 2.55)
        box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.35))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_BODY
        p2.space_before = Pt(8)

    add_footer(s9, 9)

    # =========================================================================
    # SLIDE 10: Future Roadmap & Vision 2030
    # =========================================================================
    s10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s10)
    add_header(s10, "Future Scope & Scale", "Roadmap: From School Innovation to National Digital Utility", "Structured scale-up strategy for NCSC & National Digital Education Missions", ACCENT_GOLD)

    phases = [
        ("Phase 1: School Level (2026)", 
         "Ready-to-Deploy Prototype", 
         "• 100% Free sovereign sub-user architecture deployed\n• Socratic AI engine with instant zero-cost onboarding\n• InAixVAi teacher administrative dashboard\n• Pilot launch at PM SHRI KV No. 1 AFS Chakeri", 
         ACCENT_BLUE),
        ("Phase 2: District & State Level", 
         "Vernacular & Voice Scaling", 
         "• Integration of 10+ Indian regional languages\n• Bi-directional voice conversation engine\n• Offline mesh/edge caching for remote village schools\n• State-wide KV & Government school rollout", 
         ACCENT_GREEN),
        ("Phase 3: National Level", 
         "Public Sovereign AI Ecosystem", 
         "• Potential integration with DIKSHA & PM eVidya\n• Interactive AI virtual science & coding lab assistant\n• Complete Atmanirbhar Bharat AI mentoring stack", 
         ACCENT_GOLD)
    ]

    for idx, (phase, subtitle, bullets, col) in enumerate(phases):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(1.8)
        box = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.75), Inches(5.0))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.8)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = phase
        p.font.name = "Segoe UI"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = col
        
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = "Segoe UI"
        p_sub.font.size = Pt(13)
        p_sub.font.color.rgb = TEXT_DARK
        p_sub.font.bold = True
        p_sub.space_before = Pt(4)

        p_b = tf.add_paragraph()
        p_b.text = bullets
        p_b.font.name = "Segoe UI"
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = TEXT_BODY
        p_b.space_before = Pt(12)

    add_footer(s10, 10)

    # =========================================================================
    # SLIDE 11: Conclusion & Jury Q&A Summary
    # =========================================================================
    s11 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s11)

    # Center card (Dark Executive Header Theme matching Report)
    card_end = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(0.8), Inches(10.933), Inches(5.8))
    card_end.fill.solid()
    card_end.fill.fore_color.rgb = HERO_DARK
    card_end.line.color.rgb = ACCENT_CYAN
    card_end.line.width = Pt(2)
    tf_end = card_end.text_frame
    tf_end.word_wrap = True

    p0 = tf_end.paragraphs[0]
    p0.text = "🏆 INVENTED FOR BHARAT'S FUTURE LEARNERS | NCSC 2026-27"
    p0.font.name = "Segoe UI"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_YELLOW
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf_end.add_paragraph()
    p1.text = "Thank You, Respected Jury Members!"
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(8)

    p2 = tf_end.add_paragraph()
    p2.text = "V.A.N.I - xAI (BoVxAi Edition) : Vāṇī Adhyātmik Navīn Intellect"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(19)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_MINT
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)

    p3 = tf_end.add_paragraph()
    p3.text = "\"Democratizing cutting-edge AI mentorship for every student in every corner of India — 100% Free, Zero Subscriptions, Zero E-Waste & Grounded in Scientific Integrity.\""
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(13.5)
    p3.font.color.rgb = TEXT_WHITE
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(12)

    # 4 Summary Highlight Chips
    chips_text = "✨ 100% Free (Zero Subscriptions)  |  🔑 Algorithmic UIDs  |  🎓 Socratic Mentoring  |  💰 ₹10,000 DBT Scalability"
    p_chip = tf_end.add_paragraph()
    p_chip.text = chips_text
    p_chip.font.name = "Segoe UI"
    p_chip.font.size = Pt(12)
    p_chip.font.bold = True
    p_chip.font.color.rgb = ACCENT_CYAN
    p_chip.alignment = PP_ALIGN.CENTER
    p_chip.space_before = Pt(12)

    p4 = tf_end.add_paragraph()
    p4.text = "Innovator: Dhruv Sagar (Class 10th) | PM SHRI KV NO.1 AFS Chakeri, Kanpur (UP)\nNational Children's Science Congress (NCSC) | NCSTC, DST India"
    p4.font.name = "Segoe UI"
    p4.font.size = Pt(12.5)
    p4.font.color.rgb = RGBColor(203, 213, 225)
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(16)

    add_footer(s11, 11)

    # Save presentation
    output_filename = "VANI-xAi_IM_BoVxAi.pptx"
    prs.save(output_filename)
    print(f"[SUCCESS] Presentation regenerated successfully as: {output_filename}")

if __name__ == "__main__":
    create_vani_presentation()
