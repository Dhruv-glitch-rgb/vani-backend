import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Theme Colors
    BG_DARK = RGBColor(15, 23, 42)       # #0F172A (Deep Slate)
    CARD_BG = RGBColor(30, 41, 59)       # #1E293B (Card Dark)
    CARD_BORDER = RGBColor(51, 65, 85)   # #334155 (Subtle border)
    ACCENT_CYAN = RGBColor(56, 189, 248) # #38BDF8 (Vibrant Cyan)
    ACCENT_GREEN = RGBColor(52, 211, 153)# #34D399 (Emerald Green)
    ACCENT_GOLD = RGBColor(251, 191, 36) # #FBBF24 (Amber/Gold)
    TEXT_WHITE = RGBColor(248, 250, 252) # #F8FAFC (Pure White)
    TEXT_MUTED = RGBColor(148, 163, 184) # #94A3B8 (Light Slate)
    TEXT_LIGHT = RGBColor(226, 232, 240) # #E2E8F0 (Off White)

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return bg

    def add_header(slide, category_tag, title_text, subtitle_text=None):
        # Category Tag badge
        tag_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.45), Inches(3.4), Inches(0.35))
        tag_box.fill.solid()
        tag_box.fill.fore_color.rgb = CARD_BG
        tag_box.line.color.rgb = ACCENT_CYAN
        tag_box.line.width = Pt(1)
        tf = tag_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_tag.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        p.alignment = PP_ALIGN.CENTER

        # Title
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.7))
        tf_t = tb.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(26)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE

        if subtitle_text:
            p_s = tf_t.add_paragraph()
            p_s.text = subtitle_text
            p_s.font.size = Pt(13)
            p_s.font.color.rgb = TEXT_MUTED
            p_s.space_before = Pt(4)

    # ----------------------------------------------------
    # SLIDE 1: Title Slide
    # ----------------------------------------------------
    s1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s1)

    # Top Pill
    pill = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.3), Inches(0.7), Inches(4.7), Inches(0.45))
    pill.fill.solid()
    pill.fill.fore_color.rgb = CARD_BG
    pill.line.color.rgb = ACCENT_GOLD
    pill.line.width = Pt(1.5)
    tf_p = pill.text_frame
    p0 = tf_p.paragraphs[0]
    p0.text = "🏆 NATIONAL CHILDREN'S SCIENCE CONGRESS (NCSC) 2026-27"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_GOLD
    p0.alignment = PP_ALIGN.CENTER

    # Main Title Box
    title_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.4), Inches(11.333), Inches(2.2))
    tf_main = title_box.text_frame
    tf_main.word_wrap = True
    p1 = tf_main.paragraphs[0]
    p1.text = "V.A.N.I - xAI"
    p1.font.size = Pt(54)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_CYAN
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf_main.add_paragraph()
    p2.text = "Vāṇī Adhyātmik Navīn Intellect"
    p2.font.size = Pt(22)
    p2.font.color.rgb = ACCENT_GREEN
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)

    p3 = tf_main.add_paragraph()
    p3.text = "Next-Generation 100% Free Sovereign AI Tutoring & Innovation Ecosystem for Indian Students"
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_LIGHT
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(8)

    # Info Cards Container (Student & Guide Info)
    card1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(3.9), Inches(5.2), Inches(2.8))
    card1.fill.solid()
    card1.fill.fore_color.rgb = CARD_BG
    card1.line.color.rgb = CARD_BORDER
    tf_c1 = card1.text_frame
    tf_c1.word_wrap = True
    c1_p1 = tf_c1.paragraphs[0]
    c1_p1.text = "👤 Innovator Details"
    c1_p1.font.size = Pt(16)
    c1_p1.font.bold = True
    c1_p1.font.color.rgb = ACCENT_CYAN

    items_c1 = [
        ("Student Name:", "Dhruv Sagar"),
        ("Class & Stream:", "Class 10th (Secondary)"),
        ("Institution:", "PM SHRI KV NO.1 AFS Chakeri, Kanpur"),
        ("State / Region:", "Uttar Pradesh, India")
    ]
    for k, v in items_c1:
        p = tf_c1.add_paragraph()
        p.text = f"• {k} {v}"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(6)

    card2 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(3.9), Inches(5.2), Inches(2.8))
    card2.fill.solid()
    card2.fill.fore_color.rgb = CARD_BG
    card2.line.color.rgb = CARD_BORDER
    tf_c2 = card2.text_frame
    tf_c2.word_wrap = True
    c2_p1 = tf_c2.paragraphs[0]
    c2_p1.text = "🎯 Project Attributes"
    c2_p1.font.size = Pt(16)
    c2_p1.font.bold = True
    c2_p1.font.color.rgb = ACCENT_GREEN

    items_c2 = [
        ("Project Type:", "Applied AI / Educational Technology"),
        ("Accessibility:", "100% Free Lifetime (Zero Subscriptions)"),
        ("Architecture:", "Asymmetric Cloud-Native Sub-User (<50KB)"),
        ("Hardware:", "100% Zero E-Waste (Software Solution)")
    ]
    for k, v in items_c2:
        p = tf_c2.add_paragraph()
        p.text = f"• {k} {v}"
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(6)

    # ----------------------------------------------------
    # SLIDE 2: Grassroots Problem Statement
    # ----------------------------------------------------
    s2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s2)
    add_header(s2, "Problem Statement", "The Digital Learning Divide in India", "Key challenges faced by millions of school students in Tier-2, Tier-3 and rural areas")

    problems = [
        ("💸 Expensive AI Paywalls", "Commercial tools like ChatGPT Plus cost ₹1,500-₹2,500/month, making advanced AI mentoring inaccessible to middle and low-income students."),
        ("📱 Hardware Barriers", "State-of-the-art AI apps require high-end smartphones and heavy RAM, excluding students with basic devices and 2G/3G connectivity."),
        ("🧠 The 'Answer Machine' Trap", "Standard AI regurgitates direct answers, encouraging cheating and passive copying rather than actual cognitive understanding."),
        ("🛡️ Lack of Localized Safety", "Mainstream models lack Indian curriculum contextualization, vernacular safety filters, and ethical guardrails designed for school children.")
    ]

    for idx, (head, desc) in enumerate(problems):
        row = idx // 2
        col_idx = idx % 2
        left = Inches(0.8 + col_idx * 5.95)
        top = Inches(1.8 + row * 2.55)
        box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.35))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = CARD_BORDER
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = head
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(248, 113, 113) # Soft Red
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_LIGHT
        p2.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 3: The Innovation - V.A.N.I-xAI Solution
    # ----------------------------------------------------
    s3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s3)
    add_header(s3, "The Proposed Innovation", "V.A.N.I-xAI: Sovereign AI Mentor for Every Student", "Democratizing cloud intelligence through lightweight asymmetric architecture")

    solutions = [
        ("⚡ Asymmetric Sub-User Core", "99% of computation happens on optimized cloud nodes. The client is ultralight (<50KB), loading in under 1 second on low-end phones & 2G connections."),
        ("🎓 Socratic Mentorship Engine", "Engineered to act as a digital mentor—guiding students with step-by-step hints and concept scaffolding rather than feeding raw answers."),
        ("🔓 100% Free Sovereign Access", "Removed all premium subscription plans and paywalls. Every capability is 100% free and permanently unlocked for all students."),
        ("🌐 Hyper-Localized Indian Context", "Built with custom system prompts reflecting Indian curricula (CBSE/State Boards), regional analogies, and strict youth safety protocols.")
    ]

    for idx, (title, desc) in enumerate(solutions):
        row = idx // 2
        col_idx = idx % 2
        left = Inches(0.8 + col_idx * 5.95)
        top = Inches(1.8 + row * 2.55)
        box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.35))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = CARD_BORDER
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"0{idx+1}. {title}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_LIGHT
        p2.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 4: Technological Uniqueness & Core Inventions
    # ----------------------------------------------------
    s4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s4)
    add_header(s4, "Scientific Novelty & IP", "4 Pillars of Novel Technological Architecture", "Proprietary features developed from scratch for the V.A.N.I ecosystem")

    pillars = [
        ("1. Mathematical UID Generation", "Generates permanent algorithmic IDs (e.g., V.A.N.I-xAI-DHRU-3478) upon first OAuth login, binding persistent workspace state directly to hardware-agnostic credentials.", ACCENT_CYAN),
        ("2. 3-Tier Cryptographic Security", "Algorithmic authorization validated against 3 cryptographic checks (DB Existence, Active Token, and Strict User UID Match) preventing identity spoofing.", ACCENT_GREEN),
        ("3. InAixVAi Command Dashboard", "Dedicated educator/admin portal offering real-time telemetry, concurrent sub-user monitoring, query management, and security oversight with zero subscription barriers.", ACCENT_GOLD),
        ("4. Autonomous Session & Memory Protocol", "Autonomous self-healing lifecycle engine that continuously manages SQLite memory persistence (`vani_memory.db`), multi-device sync, and intelligent load-balancing.", RGBColor(192, 132, 252))
    ]

    for idx, (title, desc, col) in enumerate(pillars):
        left = Inches(0.8 + idx * 2.95)
        top = Inches(1.8)
        box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(5.1))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_LIGHT
        p2.space_before = Pt(12)

    # ----------------------------------------------------
    # SLIDE 5: End-to-End System Architecture & Workflow
    # ----------------------------------------------------
    s5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s5)
    add_header(s5, "System Architecture", "5-Step Operational Flow & Data Pipeline", "Seamless, secure, and zero-cost interaction from login to AI inference")

    steps = [
        ("Step 1", "Auth & ID Gen", "Student signs in via Google OAuth. System algorithmically generates user-bound V.A.N.I UID."),
        ("Step 2", "Instant Free Access", "Student launches workspace with 100% unlocked features, zero fees, and instant access."),
        ("Step 3", "Security & Swarm Sync", "System cryptographically binds UID with session tokens for cross-device sync and memory persistence."),
        ("Step 4", "Socratic AI Inference", "Student queries AI; server formats prompt with Socratic scaffolds & guardrails in real-time."),
        ("Step 5", "Continuous Learning", "Autonomous memory engine records learning progress, concepts mastered, and voice sessions.")
    ]

    for idx, (step_num, title, desc) in enumerate(steps):
        left = Inches(0.8 + idx * 2.38)
        top = Inches(2.0)
        box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.25), Inches(4.7))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = ACCENT_CYAN if idx % 2 == 0 else ACCENT_GREEN
        box.line.width = Pt(1.5)
        
        tf = box.text_frame
        tf.word_wrap = True
        
        # Step number badge
        p0 = tf.paragraphs[0]
        p0.text = step_num.upper()
        p0.font.size = Pt(12)
        p0.font.bold = True
        p0.font.color.rgb = ACCENT_GOLD
        
        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE
        p1.space_before = Pt(6)
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(10)

    # ----------------------------------------------------
    # SLIDE 6: Comparative Advantage Matrix
    # ----------------------------------------------------
    s6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s6)
    add_header(s6, "Competitive Benchmarking", "V.A.N.I-xAI vs. Conventional AI & EdTech", "Engineered for maximum accessibility and educational integrity")

    # Table creation
    rows = 6
    cols = 4
    left = Inches(0.8)
    top = Inches(1.8)
    width = Inches(11.733)
    height = Inches(4.9)

    table_shape = s6.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(2.8)
    table.columns[2].width = Inches(2.8)
    table.columns[3].width = Inches(2.933)

    headers = ["Evaluation Metric", "Commercial AI (ChatGPT/Claude)", "Traditional EdTech Apps", "V.A.N.I - xAI (Our Solution)"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = ACCENT_GOLD if i == 3 else ACCENT_CYAN

    data = [
        ("Cost for Students", "₹1,500 - ₹2,500/month (Expensive)", "₹5,000 - ₹20,000/yr (Commercial)", "100% Free Lifetime (Zero Subscriptions)"),
        ("Device Hardware Barrier", "Requires modern smartphone & RAM", "Requires 100MB+ app download", "Runs on any browser / <50KB payload"),
        ("Educational Approach", "Direct answer generation (cheat)", "Static video lectures (Passive)", "Socratic guidance & Concept mentor"),
        ("Payment & Access Barrier", "Requires International Cards", "Closed recurring subscriptions", "Zero Paywalls / Instant Free Access"),
        ("Zero E-Waste & Ecology", "Heavy local battery drain", "Hardware locks / tablets", "100% Cloud-Native (Zero E-waste)")
    ]

    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(20, 30, 48) if r_idx % 2 == 0 else CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(11)
            if c_idx == 3:
                p.font.bold = True
                p.font.color.rgb = ACCENT_GREEN
            elif c_idx == 0:
                p.font.bold = True
                p.font.color.rgb = TEXT_WHITE
            else:
                p.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 7: Social, Environmental Impact & DBT ₹10,000 Utilization
    # ----------------------------------------------------
    s7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s7)
    add_header(s7, "Grassroots Impact & Scaling Plan", "Social Value, Sustainability & Scalability", "Maximizing grassroots educational impact for Bharat")

    # 3 Large Cards
    c1 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5.1))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = ACCENT_CYAN
    c1.line.width = Pt(1.5)
    tf1 = c1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "🌱 100% Zero E-Waste"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    
    points1 = [
        "No toxic batteries, PCB boards, or plastic chassis manufacturing.",
        "Saves thousands of kilograms of electronic waste compared to hardware gadgets.",
        "Repurposes existing school computers and parents' low-end smartphones.",
        "Eco-friendly, energy-efficient server inference."
    ]
    for pt in points1:
        p = tf1.add_paragraph()
        p.text = f"• {pt}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(8)

    c2 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5.1))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = ACCENT_GREEN
    c2.line.width = Pt(1.5)
    tf2 = c2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "🤝 Grassroots Empowerment"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    
    points2 = [
        "Brings 24/7 personal tutor access to students in rural & government schools with zero cost.",
        "Assists first-generation learners in science, math, coding, and language arts.",
        "Prevents digital divide from widening between metro and rural students.",
        "Sovereign Indian innovation aligned with NEP 2020 & Digital India vision."
    ]
    for pt in points2:
        p = tf2.add_paragraph()
        p.text = f"• {pt}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(8)

    c3 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.1))
    c3.fill.solid()
    c3.fill.fore_color.rgb = CARD_BG
    c3.line.color.rgb = ACCENT_GOLD
    c3.line.width = Pt(1.5)
    tf3 = c3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "💰 ₹10,000 DBT Fund Plan"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    
    points3 = [
        "₹5,000: High-speed Cloud API credits (empowers 100,000+ student queries at zero cost).",
        "₹3,000: Serverless hosting, Firebase DB & low-latency Indian edge CDN scaling.",
        "₹1,500: Vernacular Language token optimization for Hindi & Regional dialects.",
        "₹500: SSL Security certificates and automated domain hosting."
    ]
    for pt in points3:
        p = tf3.add_paragraph()
        p.text = f"• {pt}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_LIGHT
        p.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 8: Live Prototype Demonstration & Features
    # ----------------------------------------------------
    s8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s8)
    add_header(s8, "Working Prototype & Implementation", "Ready-to-Deploy Functional Capabilities", "Proven software architecture already active and operational")

    features = [
        ("🔐 InAixVAi Command Dashboard", "Real-time student monitoring, classroom session management, and query load oversight.", ACCENT_CYAN),
        ("🧠 Real-time LLM Router & Memory", "Integrated semantic search, conversational context retention (`vani_memory.db`), and autonomous agent loop.", ACCENT_GREEN),
        ("🎙️ Multimodal Voice Engine", "Integrated TTS (`voice_agent.py`) allowing audio-based learning for students with reading difficulties.", ACCENT_GOLD),
        ("📱 Cross-Platform Responsive Web UI", "Zero installation needed; accessible instantly from any standard browser on Android, iOS, Windows, or Linux.", RGBColor(244, 114, 182))
    ]

    for idx, (title, desc, col) in enumerate(features):
        row = idx // 2
        col_idx = idx % 2
        left = Inches(0.8 + col_idx * 5.95)
        top = Inches(1.8 + row * 2.55)
        box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.35))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = col
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = TEXT_LIGHT
        p2.space_before = Pt(8)

    # ----------------------------------------------------
    # SLIDE 9: Future Roadmap & Vision 2030
    # ----------------------------------------------------
    s9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s9)
    add_header(s9, "Future Scope & Scale", "Scaling V.A.N.I-xAI Across Indian Classrooms", "Roadmap from school-level prototype to nationwide public digital utility")

    phases = [
        ("Phase 1: Present (2026)", "Working Prototype", "• 100% Free sovereign sub-user architecture deployed\n• Socratic AI engine with instant zero-cost access\n• InAixVAi administrative command center\n• Deployed live on cloud", ACCENT_CYAN),
        ("Phase 2: District Level", "Vernacular Expansion", "• 100+ Regional Indian languages (Hindi, Tamil, Bengali, Marathi, etc.)\n• Voice-to-voice bidirectional AI tutor\n• KV / PM SHRI school pilot program", ACCENT_GREEN),
        ("Phase 3: National Scale", "Public Good Ecosystem", "• Integration with DIKSHA & PM eVidya\n• Offline mesh/edge caching for zero-internet zones\n• AI coding & scientific experiment lab partner", ACCENT_GOLD)
    ]

    for idx, (phase, subtitle, bullets, col) in enumerate(phases):
        left = Inches(0.8 + idx * 3.95)
        top = Inches(1.8)
        box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.75), Inches(5.1))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = col
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = col
        
        p_sub = tf.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.size = Pt(13)
        p_sub.font.color.rgb = TEXT_WHITE
        p_sub.font.bold = True
        p_sub.space_before = Pt(4)

        p_b = tf.add_paragraph()
        p_b.text = bullets
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = TEXT_LIGHT
        p_b.space_before = Pt(12)

    # ----------------------------------------------------
    # SLIDE 10: Conclusion & Jury Q&A
    # ----------------------------------------------------
    s10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s10)

    # Center card
    card_end = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.0), Inches(10.333), Inches(5.5))
    card_end.fill.solid()
    card_end.fill.fore_color.rgb = CARD_BG
    card_end.line.color.rgb = ACCENT_CYAN
    card_end.line.width = Pt(2)
    tf_end = card_end.text_frame
    tf_end.word_wrap = True

    p0 = tf_end.paragraphs[0]
    p0.text = "🏆 INVENTED FOR BHARAT'S FUTURE LEARNERS"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_GOLD
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf_end.add_paragraph()
    p1.text = "Thank You, Respected Jury Members!"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(10)

    p2 = tf_end.add_paragraph()
    p2.text = "V.A.N.I-xAI: Vāṇī Adhyātmik Navīn Intellect"
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_GREEN
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)

    p3 = tf_end.add_paragraph()
    p3.text = "\"Democratizing cutting-edge AI for every student in every corner of India — 100% Free, Zero Subscriptions, Zero E-Waste & Complete Scientific Integrity.\""
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_LIGHT
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(14)

    p4 = tf_end.add_paragraph()
    p4.text = "Innovator: Dhruv Sagar (Class 10th) | PM SHRI KV NO.1 AFS Chakeri Kanpur\nNational Children's Science Congress (NCSC) | NCSTC, DST India"
    p4.font.size = Pt(13)
    p4.font.color.rgb = TEXT_MUTED
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(20)

    # Save presentation
    output_path = "VANI_xAI_NCSC_2026.pptx"
    prs.save(output_path)
    print(f"Presentation generated successfully at: {output_path}")

if __name__ == "__main__":
    create_presentation()
